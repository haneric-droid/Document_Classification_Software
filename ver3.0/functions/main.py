import csv
import re
import time
import zipfile
from collections import Counter
from io import BytesIO, StringIO
from typing import Any, List
from xml.etree import ElementTree

import fitz
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from firebase_functions import https_fn
from firebase_functions.params import StringParam
from firebase_admin import initialize_app

initialize_app()

MAX_EVIDENCE_CHARS = 4500
MAX_KEYWORD_SENTENCES = 30
TEXT_LAYER_MIN_CHARS = 50

KEYWORDS = [
    "공고", "모집", "안내", "신청서", "제출서류", "접수기간", "양식", "서식", "붙임", "별첨",
    "사업계획", "수행계획", "추진전략", "연구목표", "개발목표", "예산", "일정", "기대효과", "제안서",
    "조사", "설문", "인터뷰", "통계", "시장분석", "사례", "리서치", "응답자",
    "발표", "슬라이드", "PPT", "발표자료", "최종발표", "중간발표",
    "계약서", "갑", "을", "대금", "지급", "정산", "계좌", "입금", "송금", "거래 조건",
    "보고서", "리포트", "고찰", "결론", "요약", "서론", "본론", "참고문헌", "시사점", "배경", "연구",
    "회의록", "회의", "안건", "참석자", "논의", "결정사항", "액션아이템", "메모", "기록", "공유사항",
    "교육", "학습", "매뉴얼", "가이드", "사용법", "온보딩", "문제", "정답", "해설", "연습문제", "과제", "시험",
    "채용", "입사", "면접", "이력서", "자기소개서", "근로계약", "휴가", "인사평가", "급여", "퇴사",
    "매출", "비용", "지출", "영수증", "세금계산서", "재무제표", "손익", "회계", "결산",
    "구매", "발주", "견적", "납품", "입고", "거래명세서", "공급업체", "단가", "수량",
    "고객", "영업", "상담", "문의", "미팅", "니즈", "요구사항", "거래처", "CRM",
    "마케팅", "홍보", "광고", "캠페인", "콘텐츠", "SNS", "브랜드", "보도자료",
    "API", "시스템", "설계", "데이터베이스", "ERD", "테스트", "배포", "서버",
    "증명서", "인증서", "자격증", "수료증", "사업자등록증", "재직증명서",
]

KEY_SENTENCE_PRIORITY_TERMS = [
    "공고", "신청", "제출", "계약", "정산", "조사", "분석", "발표", "계획", "예산", "일정", "목표",
    "모집", "접수", "지급", "청구", "세금계산서", "설문", "인터뷰", "제안", "수행",
    "보고서", "회의록", "안건", "결정사항", "교육", "매뉴얼", "채용", "면접",
    "재무제표", "손익", "회계", "구매", "발주", "납품", "고객", "영업",
    "API", "요구사항", "시스템", "설계", "테스트", "배포",
    "증명서", "인증서", "자격증", "발급기관", "등록번호",
]

STOP_WORDS = {
    "그리고", "그러나", "입니다", "합니다", "있는", "없는", "대한", "관련", "문서", "자료",
    "파일", "페이지", "내용", "확인", "제출", "첨부", "기반", "경우", "또는", "위한",
    "the", "and", "for", "with", "from", "this", "that",
}


# ── 유틸 ──────────────────────────────────────────────────────────────
def normalize_whitespace(text: str) -> str:
    return re.sub(r"\s+", " ", text or "").strip()

def truncate_text(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    return text[:max_chars - 20].rstrip() + "\n... [truncated]"

def decode_text_bytes(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "cp949", "euc-kr"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="ignore")

def extract_xml_text(data: bytes) -> str:
    try:
        root = ElementTree.fromstring(data)
        return " ".join(t.strip() for t in root.itertext() if t and t.strip())
    except Exception:
        raw = decode_text_bytes(data)
        return normalize_whitespace(re.sub(r"<[^>]+>", " ", raw))

def find_keyword_sentences(text: str) -> List[str]:
    normalized = normalize_whitespace(text)
    sentences = re.split(r"(?<=[.!?。！？])\s+|\n+", normalized)
    found: List[str] = []
    seen: set = set()
    for sentence in sentences:
        clean = sentence.strip()
        if len(clean) < 8:
            continue
        if any(kw.lower() in clean.lower() for kw in KEYWORDS) or \
           any(t.lower() in clean.lower() for t in KEY_SENTENCE_PRIORITY_TERMS):
            shortened = truncate_text(clean, 220)
            key = shortened.lower()
            if key not in seen:
                seen.add(key)
                found.append(shortened)
        if len(found) >= MAX_KEYWORD_SENTENCES:
            break
    return found

def extract_key_sentence_candidates(text: str) -> List[str]:
    normalized = normalize_whitespace(text)
    sentences = re.split(r"(?<=[.!?。！？])\s+|\n+|(?<=다\.)\s+", normalized)
    scored: list = []
    seen: set = set()
    date_p = re.compile(r"\d{4}[./-]\d{1,2}[./-]\d{1,2}")
    money_p = re.compile(r"\d[\d,]*(원|만원|억원|KRW|USD)", re.IGNORECASE)
    for sentence in sentences:
        clean = sentence.strip()
        if len(clean) < 18:
            continue
        shortened = truncate_text(clean, 240)
        key = shortened.lower()
        if key in seen:
            continue
        seen.add(key)
        score = 0
        lower = clean.lower()
        if any(t.lower() in lower for t in KEY_SENTENCE_PRIORITY_TERMS):
            score += 4
        if date_p.search(clean):
            score += 3
        if money_p.search(clean):
            score += 3
        if score:
            scored.append((score, shortened))
    scored.sort(key=lambda x: (-x[0], len(x[1])))
    return [s for _, s in scored[:10]]

def top_frequent_terms(text: str) -> List[tuple]:
    tokens = re.findall(r"[가-힣A-Za-z0-9]{2,}", text.lower())
    filtered = [t for t in tokens if t not in STOP_WORDS and not t.isdigit()]
    return Counter(filtered).most_common(20)

def build_text_evidence_package(filename: str, text_value: str, source_label: str, page_count: int = 0, count_label: str = "총 페이지 수") -> str:
    normalized = normalize_whitespace(text_value)
    first_text = truncate_text(normalized[:1800], 1800)
    last_text = truncate_text(normalized[-1200:], 1200) if len(normalized) > 1200 else normalized
    keyword_sentences = find_keyword_sentences(normalized)
    top_words = top_frequent_terms(normalized)
    key_sentences = extract_key_sentence_candidates(normalized)
    sections = [
        f"파일명: {filename}",
        f"추출 방식: {source_label}",
        f"전체 텍스트 길이: {len(normalized)}",
    ]
    if page_count:
        sections.append(f"{count_label}: {page_count}")
    sections.extend([
        "[첫 부분 텍스트]", first_text or "추출된 텍스트가 없습니다.",
        "[마지막 부분 텍스트]", last_text or "추출된 텍스트가 없습니다.",
        "[분류 키워드 포함 문장 최대 30개]",
        "\n".join(f"- {s}" for s in keyword_sentences) or "없음",
        "[자주 등장하는 단어 Top 20]",
        ", ".join(f"{w}({c})" for w, c in top_words) or "없음",
        "[핵심 문장 후보 최대 10개]",
        "\n".join(f"- {s}" for s in key_sentences) or "없음",
    ])
    return truncate_text("\n\n".join(sections), MAX_EVIDENCE_CHARS)


# ── 형식별 추출 ────────────────────────────────────────────────────────
def extract_pdf(filename: str, content: bytes) -> dict:
    try:
        with fitz.open(stream=content, filetype="pdf") as doc:
            page_count = doc.page_count
            all_text = "\n".join(normalize_whitespace(p.get_text("text")) for p in doc)
        total_len = len(normalize_whitespace(all_text))
        if total_len < TEXT_LAYER_MIN_CHARS:
            return {"extraction_status": "no_text_layer", "mode": "PDF 열기 성공 - 텍스트 레이어 없음",
                    "page_count": page_count, "evidence_package": "",
                    "warning": "스캔본 또는 이미지형 PDF일 수 있습니다. OCR 처리가 필요합니다."}
        return {"extraction_status": "success", "mode": "PyMuPDF 실제 텍스트 추출",
                "page_count": page_count,
                "evidence_package": build_text_evidence_package(filename, all_text, "PyMuPDF", page_count),
                "warning": ""}
    except Exception as e:
        return {"extraction_status": "failed", "mode": "fallback", "page_count": 0,
                "evidence_package": "", "warning": str(e)}

def extract_docx(filename: str, content: bytes) -> dict:
    try:
        doc = DocxDocument(BytesIO(content))
        chunks = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    chunks.append(" | ".join(cells))
        text_value = "\n".join(chunks)
        if len(normalize_whitespace(text_value)) < 1:
            raise ValueError("empty")
        return {"extraction_status": "success", "mode": "python-docx 실제 텍스트 추출",
                "page_count": 0,
                "evidence_package": build_text_evidence_package(filename, text_value, "python-docx"),
                "warning": ""}
    except Exception as e:
        return {"extraction_status": "failed", "mode": "fallback", "page_count": 0,
                "evidence_package": "", "warning": str(e)}

def extract_txt(filename: str, content: bytes) -> dict:
    for enc in ("utf-8-sig", "utf-8", "cp949", "euc-kr"):
        try:
            text_value = content.decode(enc)
            if not normalize_whitespace(text_value):
                continue
            return {"extraction_status": "success", "mode": f"TXT 텍스트 추출 ({enc})",
                    "page_count": 0,
                    "evidence_package": build_text_evidence_package(filename, text_value, f"TXT ({enc})"),
                    "warning": ""}
        except UnicodeDecodeError:
            continue
    return {"extraction_status": "failed", "mode": "fallback", "page_count": 0,
            "evidence_package": "", "warning": "TXT 디코딩 실패"}

def extract_hwpx(filename: str, content: bytes) -> dict:
    try:
        chunks: List[str] = []
        with zipfile.ZipFile(BytesIO(content)) as archive:
            for name in archive.namelist():
                if name.lower() == "preview/prvtext.txt" or \
                   (name.lower().startswith(("contents/", "sections/")) and name.lower().endswith(".xml")):
                    data = archive.read(name)
                    chunks.append(decode_text_bytes(data) if name.endswith(".txt") else extract_xml_text(data))
        text_value = "\n".join(c for c in chunks if c.strip())
        if len(normalize_whitespace(text_value)) < TEXT_LAYER_MIN_CHARS:
            raise ValueError("too short")
        return {"extraction_status": "success", "mode": "HWPX XML 추출",
                "page_count": 0,
                "evidence_package": build_text_evidence_package(filename, text_value, "HWPX XML"),
                "warning": ""}
    except Exception as e:
        return {"extraction_status": "failed", "mode": "fallback", "page_count": 0,
                "evidence_package": "", "warning": str(e)}

def extract_pptx_shape_text(shape) -> List[str]:
    chunks: List[str] = []
    try:
        if getattr(shape, "has_text_frame", False) and shape.text_frame:
            for para in shape.text_frame.paragraphs:
                text = " ".join(run.text for run in para.runs if run.text.strip())
                if text.strip():
                    chunks.append(text.strip())
    except Exception:
        pass
    try:
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    chunks.append(" | ".join(cells))
    except Exception:
        pass
    return chunks

def extract_pptx(filename: str, content: bytes) -> dict:
    try:
        prs = Presentation(BytesIO(content))
        slide_chunks: List[str] = []
        for i, slide in enumerate(prs.slides, 1):
            chunks: List[str] = []
            for shape in slide.shapes:
                chunks.extend(extract_pptx_shape_text(shape))
            unique = list(dict.fromkeys(c for c in chunks if c))
            if unique:
                slide_chunks.append(f"[슬라이드 {i}]\n" + "\n".join(unique))
        text_value = "\n\n".join(slide_chunks)
        slide_count = len(prs.slides)
        if len(normalize_whitespace(text_value)) < TEXT_LAYER_MIN_CHARS:
            return {"extraction_status": "no_text_layer", "mode": "PPTX - 텍스트 부족",
                    "page_count": slide_count, "evidence_package": "",
                    "warning": "이미지형 슬라이드는 OCR이 필요합니다."}
        return {"extraction_status": "success", "mode": "python-pptx 실제 텍스트 추출",
                "page_count": slide_count,
                "evidence_package": build_text_evidence_package(filename, text_value, "python-pptx", slide_count, "총 슬라이드 수"),
                "warning": ""}
    except Exception as e:
        return {"extraction_status": "failed", "mode": "fallback", "page_count": 0,
                "evidence_package": "", "warning": str(e)}

def safe_cell(value: Any) -> str:
    return normalize_whitespace(str(value))[:120] if value is not None else ""

def extract_xlsx(filename: str, content: bytes) -> dict:
    try:
        wb = load_workbook(BytesIO(content), read_only=True, data_only=True)
        sheet_names = wb.sheetnames
        rows: List[str] = []
        headers: List[str] = []
        for ws in wb.worksheets[:5]:
            first = True
            for row in ws.iter_rows(max_row=60, values_only=True):
                cells = [safe_cell(c) for c in row if safe_cell(c)]
                if not cells:
                    continue
                line = " | ".join(cells)
                if first:
                    headers.append(line)
                    first = False
                else:
                    rows.append(line)
        all_text = "\n".join(rows)
        if len(normalize_whitespace(all_text)) < TEXT_LAYER_MIN_CHARS:
            return {"extraction_status": "no_text_layer", "mode": "XLSX - 텍스트 부족",
                    "page_count": 0, "evidence_package": "", "warning": "표에서 추출 가능한 텍스트가 부족합니다."}
        sections = [f"파일명: {filename}", f"추출 방식: openpyxl",
                    f"시트명: {', '.join(sheet_names[:10])}",
                    "[헤더 후보]", "\n".join(f"- {h}" for h in headers[:10]) or "없음",
                    "[상위 행]", "\n".join(f"- {r}" for r in rows[:40]) or "없음"]
        ep = truncate_text("\n\n".join(sections), MAX_EVIDENCE_CHARS)
        return {"extraction_status": "success", "mode": "openpyxl 실제 텍스트 추출",
                "page_count": 0, "evidence_package": ep, "warning": ""}
    except Exception as e:
        return {"extraction_status": "failed", "mode": "fallback", "page_count": 0,
                "evidence_package": "", "warning": str(e)}

def extract_csv(filename: str, content: bytes) -> dict:
    for enc in ("utf-8-sig", "utf-8", "cp949", "euc-kr"):
        try:
            text = content.decode(enc)
            reader = csv.reader(StringIO(text))
            rows_raw = list(reader)
            if not rows_raw:
                continue
            headers = [" | ".join(c.strip() for c in rows_raw[0] if c.strip())]
            rows = [" | ".join(c.strip() for c in r if c.strip()) for r in rows_raw[1:51]]
            all_text = "\n".join(rows)
            if len(normalize_whitespace(all_text)) < TEXT_LAYER_MIN_CHARS:
                return {"extraction_status": "no_text_layer", "mode": "CSV - 텍스트 부족",
                        "page_count": 0, "evidence_package": "", "warning": ""}
            sections = [f"파일명: {filename}", f"추출 방식: CSV ({enc})", f"인코딩: {enc}",
                        "[헤더]", "\n".join(f"- {h}" for h in headers) or "없음",
                        "[상위 행]", "\n".join(f"- {r}" for r in rows[:40]) or "없음"]
            ep = truncate_text("\n\n".join(sections), MAX_EVIDENCE_CHARS)
            return {"extraction_status": "success", "mode": f"CSV 실제 텍스트 추출 ({enc})",
                    "page_count": 0, "evidence_package": ep, "warning": ""}
        except UnicodeDecodeError:
            continue
    return {"extraction_status": "failed", "mode": "fallback", "page_count": 0,
            "evidence_package": "", "warning": "CSV 디코딩 실패"}

def detect_file_type(filename: str, content_type: str) -> str:
    name = (filename or "").lower()
    if name.endswith(".pdf"):      return "pdf"
    if name.endswith(".docx"):     return "docx"
    if name.endswith(".txt"):      return "txt"
    if name.endswith(".hwpx"):     return "hwpx"
    if name.endswith(".hwp"):      return "hwp"
    if name.endswith(".pptx"):     return "pptx"
    if name.endswith(".xlsx"):     return "xlsx"
    if name.endswith(".csv"):      return "csv"
    if any(name.endswith(ext) for ext in (".jpg",".jpeg",".png",".gif",".webp",".bmp")):
        return "image"
    return "unknown"

def analyze_content(filename: str, file_type: str, content: bytes) -> dict:
    handlers = {
        "pdf": extract_pdf, "docx": extract_docx, "txt": extract_txt,
        "hwpx": extract_hwpx, "pptx": extract_pptx, "xlsx": extract_xlsx, "csv": extract_csv,
    }
    if file_type == "hwp":
        return {"extraction_status": "unsupported", "mode": "fallback", "page_count": 0,
                "evidence_package": "", "warning": "HWP 바이너리는 전용 파서가 필요합니다."}
    if file_type == "image":
        return {"extraction_status": "ocr_required", "mode": "OCR 필요", "page_count": 0,
                "evidence_package": "", "warning": "이미지 문서는 OCR 처리가 필요합니다."}
    handler = handlers.get(file_type)
    if not handler:
        return {"extraction_status": "unsupported", "mode": "fallback", "page_count": 0,
                "evidence_package": "", "warning": "지원하지 않는 파일 형식입니다."}
    return handler(filename, content)


# ── Cloud Function ─────────────────────────────────────────────────────
CORS_HEADERS = {
    "Access-Control-Allow-Origin": "*",
    "Access-Control-Allow-Methods": "POST, OPTIONS",
    "Access-Control-Allow-Headers": "Content-Type, Authorization",
}

@https_fn.on_request()
def analyze(req: https_fn.Request) -> https_fn.Response:
    if req.method == "OPTIONS":
        return https_fn.Response("", status=204, headers=CORS_HEADERS)
    if req.method != "POST":
        return https_fn.Response("Method not allowed", status=405, headers=CORS_HEADERS)

    start = time.perf_counter()
    try:
        if req.content_type and "multipart" in req.content_type:
            file_storage = req.files.get("file")
            if not file_storage:
                return https_fn.Response("{\"ok\":false,\"warning\":\"file 필드가 없습니다.\"}", status=400,
                                         content_type="application/json")
            filename = file_storage.filename or "uploaded_file"
            content = file_storage.read()
        else:
            return https_fn.Response("{\"ok\":false,\"warning\":\"multipart/form-data 요청이 필요합니다.\"}", status=400,
                                     content_type="application/json")

        file_type = detect_file_type(filename, "")
        result = analyze_content(filename, file_type, content)

        import json
        body = json.dumps({
            "ok": True,
            "filename": filename,
            "file_type": file_type,
            "extraction_status": result["extraction_status"],
            "mode": result["mode"],
            "page_count": result["page_count"],
            "evidence_package": result["evidence_package"],
            "processing_time_sec": round(time.perf_counter() - start, 2),
            "warning": result["warning"],
        }, ensure_ascii=False)
        return https_fn.Response(body, status=200, content_type="application/json", headers=CORS_HEADERS)

    except Exception as e:
        import json
        return https_fn.Response(
            json.dumps({"ok": False, "warning": str(e)}, ensure_ascii=False),
            status=500, content_type="application/json", headers=CORS_HEADERS
        )
