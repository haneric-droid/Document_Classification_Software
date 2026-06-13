import csv
import json
import re
import hashlib
import time
import uuid
import zipfile
from collections import Counter
from datetime import datetime
from io import BytesIO, StringIO
from pathlib import Path
from typing import Any, List
from xml.etree import ElementTree

import fitz
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from fastapi import Body, Depends, FastAPI, File, Form, HTTPException, Query, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, HTMLResponse
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from sqlalchemy import Boolean, Column, DateTime, Float, ForeignKey, Integer, String, Text, create_engine, text
from sqlalchemy.orm import Session, declarative_base, sessionmaker
import jwt
from google.oauth2 import id_token
from google.auth.transport import requests as google_requests


MAX_EVIDENCE_CHARS = 4500
MAX_KEYWORD_SENTENCES = 30
TEXT_LAYER_MIN_CHARS = 50
DUPLICATE_VERSION_THRESHOLD = 70
DEFAULT_TEAM_ID = "demo-team"
TEAM_IDS = ["demo-team", "finance-team", "project-team"]
BASE_DIR = Path(__file__).resolve().parent
DB_PATH = BASE_DIR / "dowple.db"
UPLOAD_DIR = BASE_DIR / "uploads"
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

engine = create_engine(f"sqlite:///{DB_PATH}", connect_args={"check_same_thread": False})
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
Base = declarative_base()

KEYWORDS = [
    "공고", "모집", "안내", "신청서", "제출서류", "접수기간", "양식", "서식", "붙임", "별첨",
    "사업계획", "수행계획", "추진전략", "연구목표", "개발목표", "예산", "일정", "기대효과", "제안서",
    "조사", "설문", "인터뷰", "통계", "시장분석", "사례", "리서치", "응답자",
    "발표", "슬라이드", "PPT", "발표자료", "최종발표", "중간발표",
    "계약서", "갑", "을", "대금", "지급", "정산", "계좌", "입금", "송금", "거래 조건",
    "보고서", "리포트", "고찰", "결론", "요약", "서론", "본론", "참고문헌", "시사점", "배경", "연구",
    "회의록", "회의", "안건", "참석자", "논의", "결정사항", "액션아이템", "메모", "기록", "공유사항",
    "교육", "학습", "매뉴얼", "가이드", "사용법", "온보딩", "문제", "정답", "해설", "연습문제", "과제", "시험", "알고리즘", "코드", "예제", "튜토리얼",
    "채용", "입사", "면접", "이력서", "자기소개서", "근로계약", "휴가", "인사평가", "급여", "퇴사", "인사", "지원자",
    "매출", "비용", "지출", "영수증", "세금계산서", "재무제표", "손익", "회계", "결산", "비용처리", "지출결의",
    "구매", "발주", "견적", "납품", "입고", "거래명세서", "공급업체", "단가", "수량", "구매요청", "발주서",
    "고객", "영업", "상담", "문의", "미팅", "니즈", "요구사항", "거래처", "매출기회", "CRM", "고객관리",
    "마케팅", "홍보", "광고", "캠페인", "콘텐츠", "SNS", "브랜드", "보도자료", "유입", "전환율", "프로모션",
    "API", "시스템", "설계", "데이터베이스", "ERD", "버그", "테스트", "배포", "서버", "프론트엔드", "백엔드",
    "증명서", "인증서", "자격증", "수료증", "사업자등록증", "재직증명서", "경력증명서", "납세증명서", "완납증명서",
    "인감증명서", "등기부등본", "통장사본", "면허증", "허가증", "특허증", "확인서",
    "지원 대상", "신청 방법", "제출 마감", "접수처", "문의처", "첨부파일",
    "목표", "추진 배경", "수행 방법", "일정표", "예산 계획", "기대 효과", "제안 내용",
    "응답 결과", "분석 결과", "조사 대상", "표본", "경쟁사", "트렌드", "인사이트",
    "목차", "발표 순서", "발표자", "핵심 메시지", "Q&A",
    "공급가액", "청구 금액", "지급 기한", "부가가치세",
    "문제 제기", "종합 의견", "연구 배경", "결론 도출", "참고 자료", "사례 분석",
    "회의 일시", "참석자 명단", "주요 논의", "결정 사항", "후속 조치", "담당자", "다음 회의",
    "학습 목표", "실습 문제", "예제 코드", "사용 절차", "온보딩 과정", "정답을 고르시오", "풀이 과정", "업무 절차",
    "채용 공고", "면접 평가", "입사 서류", "근로 조건", "휴가 신청", "평가 결과", "급여 내역", "퇴사 절차",
    "손익계산서", "재무상태표", "비용 내역", "예산 집행", "결산 보고", "회계 처리", "지출 승인",
    "구매 요청", "발주 번호", "납품 일정", "입고 확인", "공급업체 정보", "단가 비교", "견적 비교",
    "고객 요구사항", "상담 내역", "영업 기회", "거래처 정보", "고객 미팅", "후속 연락",
    "캠페인 성과", "광고 집행", "콘텐츠 기획", "브랜드 메시지", "보도자료 배포", "유입 분석", "전환율 분석",
    "요구사항 정의", "시스템 구조", "API 명세", "데이터베이스 설계", "테스트 케이스", "버그 리포트", "배포 절차",
    "발급일", "발급기관", "증명합니다", "위 사실을 증명함", "등록번호", "사업자등록번호", "자격번호", "유효기간", "인증번호", "대표자명",
]

KEY_SENTENCE_PRIORITY_TERMS = [
    "공고", "신청", "제출", "계약", "정산", "조사", "분석", "발표", "계획", "예산", "일정", "목표",
    "모집", "접수", "지급", "청구", "세금계산서", "시장분석", "설문", "인터뷰", "제안", "수행",
    "보고서", "리포트", "회의록", "안건", "결정사항", "교육", "매뉴얼", "가이드", "채용", "면접",
    "재무제표", "손익", "회계", "구매", "발주", "납품", "고객", "영업", "마케팅", "홍보",
    "API", "요구사항", "시스템", "설계", "데이터베이스", "테스트", "배포",
    "증명서", "인증서", "자격증", "수료증", "사업자등록증", "발급기관", "등록번호",
    "notice", "apply", "submit", "contract", "settlement", "survey", "analysis", "presentation",
    "plan", "budget", "schedule", "objective", "deadline", "email", "tel", "report", "meeting",
    "education", "manual", "finance", "purchase", "sales", "marketing", "development",
]

STOP_WORDS = {
    "그리고", "그러나", "입니다", "합니다", "있는", "없는", "대한", "관련", "문서", "자료",
    "파일", "페이지", "내용", "확인", "제출", "첨부", "기반", "경우", "또는", "위한",
    "the", "and", "for", "with", "from", "this", "that",
}

DEFAULT_FOLDER_RULES = [
    {
        "id": "notice",
        "name": "공고/양식",
        "description": "외부 기관이나 조직에서 제공한 모집 공고, 신청서, 제출서류 안내, 접수기간, 양식/서식 관련 문서",
        "keywords": "공고, 모집, 안내, 신청서, 제출서류, 접수기간, 양식, 서식, 붙임, 별첨",
        "context_terms": "지원 대상, 신청 방법, 제출 마감, 접수처, 문의처, 첨부파일",
    },
    {
        "id": "plan",
        "name": "사업계획서/수행계획서",
        "description": "사업 목표, 수행 방법, 추진 전략, 연구 목표, 예산, 일정, 기대효과를 설명하는 계획 문서",
        "keywords": "사업계획, 수행계획, 추진전략, 연구목표, 개발목표, 예산, 일정, 기대효과, 제안서",
        "context_terms": "목표, 추진 배경, 수행 방법, 일정표, 예산 계획, 기대 효과, 제안 내용",
    },
    {
        "id": "research",
        "name": "조사자료",
        "description": "설문, 인터뷰, 통계, 시장분석, 사례조사, 리서치 결과를 정리한 참고 자료",
        "keywords": "조사, 설문, 인터뷰, 통계, 시장분석, 사례, 리서치, 응답자",
        "context_terms": "응답 결과, 분석 결과, 조사 대상, 표본, 경쟁사, 트렌드, 인사이트",
    },
    {
        "id": "presentation",
        "name": "발표자료",
        "description": "발표를 목적으로 제작된 슬라이드, 중간발표, 최종발표, 발표 대본, 공유용 자료",
        "keywords": "발표, 슬라이드, PPT, 발표자료, 최종발표, 중간발표",
        "context_terms": "목차, 발표 순서, 슬라이드, 발표자, 핵심 메시지, Q&A",
    },
    {
        "id": "contract",
        "name": "계약/정산",
        "description": "계약 체결, 대금 지급, 정산, 계좌 정보, 거래 조건과 관련된 문서",
        "keywords": "계약서, 계약, 갑, 을, 대금, 지급, 정산, 계좌, 입금, 송금, 거래 조건",
        "context_terms": "계약 조건, 대금 지급, 정산 기한, 계좌 정보, 입금 확인, 송금 내역, 지급 기한",
    },
    {
        "id": "other",
        "name": "기타",
        "description": "기준이 불명확하거나 텍스트 단서가 부족해 사용자의 검토가 필요한 문서",
        "keywords": "",
        "context_terms": "임시, 참고, 미정, 기타, 분류불가",
    },
    {
        "id": "report",
        "name": "보고서/리포트",
        "description": "자유서술형 보고서, 리포트, 분석 글, 고찰문, 참고문헌이 포함된 문서",
        "keywords": "보고서, 리포트, 분석, 고찰, 결론, 요약, 서론, 본론, 참고문헌, 시사점, 배경, 연구",
        "context_terms": "문제 제기, 분석 결과, 종합 의견, 연구 배경, 결론 도출, 참고 자료, 사례 분석",
        "reason": "서론, 본론, 결론, 분석 결과 등 보고서 구조와 관련된 단서가 확인됩니다.",
    },
    {
        "id": "meeting",
        "name": "회의록/메모",
        "description": "회의 내용, 안건, 참석자, 결정사항, 후속 조치, 간단한 메모를 기록한 문서",
        "keywords": "회의록, 회의, 안건, 참석자, 논의, 결정사항, 액션아이템, 메모, 기록, 공유사항",
        "context_terms": "회의 일시, 참석자 명단, 주요 논의, 결정 사항, 후속 조치, 담당자, 다음 회의",
        "reason": "회의 안건, 참석자, 결정사항, 후속 조치 등 회의 기록 단서가 확인됩니다.",
    },
    {
        "id": "education",
        "name": "교육/학습자료",
        "description": "사내 교육자료, 업무 매뉴얼, 학습자료, 연습문제, 온보딩 자료, 사용 가이드 문서",
        "keywords": "교육, 학습, 매뉴얼, 가이드, 사용법, 온보딩, 문제, 정답, 해설, 연습문제, 과제, 시험, 알고리즘, 코드, 예제, 튜토리얼",
        "context_terms": "학습 목표, 실습 문제, 예제 코드, 사용 절차, 온보딩 과정, 정답을 고르시오, 풀이 과정, 업무 절차",
        "reason": "교육, 매뉴얼, 문제, 해설, 예제, 사용법 등 학습 목적의 단서가 확인됩니다.",
    },
    {
        "id": "hr",
        "name": "인사/채용",
        "description": "채용, 입사, 면접, 근로계약, 휴가, 인사평가, 급여, 퇴사 관련 문서",
        "keywords": "채용, 입사, 면접, 이력서, 자기소개서, 근로계약, 휴가, 인사평가, 급여, 퇴사, 인사, 지원자",
        "context_terms": "채용 공고, 면접 평가, 입사 서류, 근로 조건, 휴가 신청, 평가 결과, 급여 내역, 퇴사 절차",
        "reason": "채용, 입사, 면접, 인사평가, 근로조건 등 인사 관련 단서가 확인됩니다.",
    },
    {
        "id": "finance",
        "name": "재무/회계",
        "description": "매출, 비용, 예산, 재무제표, 손익계산서, 회계 보고, 지출 관리 문서",
        "keywords": "매출, 비용, 지출, 영수증, 세금계산서, 재무제표, 손익, 예산, 회계, 결산, 비용처리, 지출결의",
        "context_terms": "손익계산서, 재무상태표, 비용 내역, 예산 집행, 결산 보고, 회계 처리, 지출 승인",
        "reason": "매출, 비용, 예산, 회계, 결산 등 재무/회계 관련 단서가 확인됩니다.",
    },
    {
        "id": "purchase",
        "name": "구매/발주",
        "description": "구매요청, 발주, 견적, 납품, 입고, 거래명세와 관련된 문서",
        "keywords": "구매, 발주, 견적, 납품, 입고, 거래명세서, 공급업체, 단가, 수량, 구매요청, 발주서",
        "context_terms": "구매 요청, 발주 번호, 납품 일정, 입고 확인, 공급업체 정보, 단가 비교, 견적 비교",
        "reason": "구매, 발주, 견적, 납품, 입고 등 구매 프로세스 관련 단서가 확인됩니다.",
    },
    {
        "id": "sales",
        "name": "영업/고객관리",
        "description": "영업 활동, 고객 상담, 제안, 고객 요구사항, 거래처 관리 관련 문서",
        "keywords": "고객, 영업, 제안, 상담, 문의, 미팅, 니즈, 요구사항, 거래처, 매출기회, CRM, 고객관리",
        "context_terms": "고객 요구사항, 상담 내역, 영업 기회, 거래처 정보, 제안 내용, 고객 미팅, 후속 연락",
        "reason": "고객, 영업, 상담, 제안, 거래처 등 영업/고객관리 단서가 확인됩니다.",
    },
    {
        "id": "marketing",
        "name": "마케팅/홍보",
        "description": "마케팅 캠페인, 홍보자료, 광고기획, 콘텐츠, 보도자료, 브랜드 관련 문서",
        "keywords": "마케팅, 홍보, 광고, 캠페인, 콘텐츠, SNS, 브랜드, 보도자료, 유입, 전환율, 프로모션",
        "context_terms": "캠페인 성과, 광고 집행, 콘텐츠 기획, 브랜드 메시지, 보도자료 배포, 유입 분석, 전환율 분석",
        "reason": "마케팅, 홍보, 광고, 캠페인, 콘텐츠 관련 단서가 확인됩니다.",
    },
    {
        "id": "tech",
        "name": "기술/개발문서",
        "description": "API 문서, 요구사항 정의서, 시스템 설계서, 데이터베이스, 테스트, 배포 관련 문서",
        "keywords": "API, 요구사항, 시스템, 설계, 데이터베이스, ERD, 버그, 테스트, 배포, 서버, 프론트엔드, 백엔드, 코드",
        "context_terms": "요구사항 정의, 시스템 구조, API 명세, 데이터베이스 설계, 테스트 케이스, 버그 리포트, 배포 절차",
        "reason": "API, 요구사항, 시스템 설계, 테스트, 배포 등 기술/개발 관련 단서가 확인됩니다.",
    },
    {
        "id": "certificate",
        "name": "증명서/인증서",
        "description": "사업자등록증, 재직증명서, 경력증명서, 수료증, 자격증, 인허가증, 납세증명서, 각종 공식 증빙 문서",
        "keywords": "증명서, 인증서, 자격증, 수료증, 사업자등록증, 재직증명서, 경력증명서, 납세증명서, 완납증명서, 인감증명서, 등기부등본, 통장사본, 면허증, 허가증, 특허증, 확인서",
        "context_terms": "발급일, 발급기관, 증명합니다, 위 사실을 증명함, 등록번호, 사업자등록번호, 자격번호, 유효기간, 인증번호, 대표자명",
        "reason": "발급기관, 등록번호, 증명 문구, 유효기간 등 공식 증빙 문서의 단서가 확인됩니다.",
    },
]

FOLDER_ICON_MAP = {
    "notice": "ti-speakerphone",
    "plan": "ti-checkup-list",
    "research": "ti-chart-dots-3",
    "presentation": "ti-presentation",
    "contract": "ti-file-invoice",
    "other": "ti-folder-question",
    "report": "ti-report",
    "meeting": "ti-notes",
    "education": "ti-school",
    "hr": "ti-user-search",
    "finance": "ti-calculator",
    "purchase": "ti-shopping-cart",
    "sales": "ti-handshake",
    "marketing": "ti-speakerphone",
    "tech": "ti-code",
    "certificate": "ti-certificate",
}


JWT_SECRET = "dowple-dms-secret-key-change-in-production"
JWT_ALGORITHM = "HS256"
# Google OAuth Client ID (환경 변수 등으로 오버라이드할 수 있도록 빈 문자열도 처리하게 구현 예정)
GOOGLE_CLIENT_ID = "884087703847-evae2cpljosmq59e75a4ukdlgsls9mvq.apps.googleusercontent.com" 

class User(Base):
    __tablename__ = "users"

    id = Column(Integer, primary_key=True, index=True)
    email = Column(String(255), nullable=False, unique=True, index=True)
    name = Column(String(255), nullable=True)
    picture = Column(Text, nullable=True)
    created_at = Column(DateTime, nullable=False, default=datetime.utcnow)


class Team(Base):
    __tablename__ = "teams"

    id = Column(String(100), primary_key=True, index=True)
    name = Column(String(255), nullable=False)
    owner_id = Column(Integer, ForeignKey("users.id"), nullable=True)
    invite_code = Column(String(50), nullable=True, unique=True, index=True)
    created_at = Column(DateTime, nullable=False, default=datetime.utcnow)


class UserTeam(Base):
    __tablename__ = "user_teams"

    user_id = Column(Integer, ForeignKey("users.id"), primary_key=True)
    team_id = Column(String(100), ForeignKey("teams.id"), primary_key=True)
    role = Column(String(50), nullable=False, default="member")  # "owner" or "member"
    joined_at = Column(DateTime, nullable=False, default=datetime.utcnow)


class Space(Base):
    __tablename__ = "spaces"

    id = Column(String(100), primary_key=True, index=True)
    team_id = Column(String(100), ForeignKey("teams.id"), nullable=False, index=True)
    name = Column(String(255), nullable=False)
    description = Column(Text, nullable=True, default="")
    created_at = Column(DateTime, nullable=False, default=datetime.utcnow)
    updated_at = Column(DateTime, nullable=False, default=datetime.utcnow, onupdate=datetime.utcnow)


class Document(Base):
    __tablename__ = "documents"

    id = Column(Integer, primary_key=True, index=True)
    team_id = Column(String(100), nullable=True, default=DEFAULT_TEAM_ID, index=True)
    space_id = Column(String(100), ForeignKey("spaces.id"), nullable=True, index=True)
    filename = Column(String(255), nullable=False)
    file_type = Column(String(50), nullable=False, default="document")
    category = Column(String(255), nullable=True, default="")
    category_id = Column(String(100), nullable=True, default="")
    summary = Column(Text, nullable=True, default="")
    confidence = Column(Float, nullable=True, default=0)
    evidence_package = Column(Text, nullable=True, default="")
    extraction_status = Column(String(100), nullable=False, default="unsupported")
    page_count = Column(Integer, nullable=False, default=0)
    file_hash = Column(String(64), nullable=True, index=True)
    version_of_id = Column(Integer, nullable=True)
    stored_path = Column(Text, nullable=False)
    created_at = Column(DateTime, nullable=False, default=datetime.utcnow)
    updated_at = Column(DateTime, nullable=False, default=datetime.utcnow, onupdate=datetime.utcnow)
    archived = Column(Boolean, nullable=False, default=False)


class FolderRule(Base):
    __tablename__ = "folder_rules"

    id = Column(String(100), primary_key=True, index=True)
    team_id = Column(String(100), nullable=True, default=DEFAULT_TEAM_ID, index=True)
    space_id = Column(String(100), ForeignKey("spaces.id"), nullable=True, index=True)
    folder_id = Column(String(100), nullable=True, index=True)
    parent_folder_id = Column(String(100), nullable=True, index=True)  # 같은 스페이스 내 부모 폴더의 folder_id (없으면 최상위)
    name = Column(String(255), nullable=False)
    icon = Column(String(100), nullable=True, default="")
    description = Column(Text, nullable=True, default="")
    keywords = Column(Text, nullable=True, default="")
    context_terms = Column(Text, nullable=True, default="")
    reason = Column(Text, nullable=True, default="")
    created_at = Column(DateTime, nullable=False, default=datetime.utcnow)
    updated_at = Column(DateTime, nullable=False, default=datetime.utcnow, onupdate=datetime.utcnow)


class ReviewQueue(Base):
    __tablename__ = "review_queue"

    id = Column(Integer, primary_key=True, index=True)
    team_id = Column(String(100), nullable=True, default=DEFAULT_TEAM_ID, index=True)
    space_id = Column(String(100), ForeignKey("spaces.id"), nullable=True, index=True)
    document_id = Column(Integer, ForeignKey("documents.id"), nullable=True)
    filename = Column(String(255), nullable=False)
    reason = Column(Text, nullable=True, default="")
    suggested_category = Column(String(255), nullable=True, default="")
    suggested_category_id = Column(String(100), nullable=True, default="")
    confidence = Column(Float, nullable=True, default=0)
    evidence_package = Column(Text, nullable=True, default="")
    recommendations_json = Column(Text, nullable=True, default="[]")
    status = Column(String(50), nullable=False, default="pending")
    resolved_category = Column(String(255), nullable=True, default="")
    resolved_category_id = Column(String(100), nullable=True, default="")
    resolved_at = Column(DateTime, nullable=True)
    created_at = Column(DateTime, nullable=False, default=datetime.utcnow)
    updated_at = Column(DateTime, nullable=False, default=datetime.utcnow, onupdate=datetime.utcnow)


security_bearer = HTTPBearer(auto_error=False)

def create_jwt_token(user_id: int, email: str) -> str:
    payload = {
        "sub": str(user_id),
        "email": email,
        "exp": time.time() + (60 * 60 * 24 * 7)  # 7일 유효
    }
    return jwt.encode(payload, JWT_SECRET, algorithm=JWT_ALGORITHM)

def decode_jwt_token(token: str) -> dict | None:
    try:
        return jwt.decode(token, JWT_SECRET, algorithms=[JWT_ALGORITHM])
    except Exception:
        return None


def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()


def get_current_user(credentials: HTTPAuthorizationCredentials = Depends(security_bearer), db: Session = Depends(get_db)) -> User:
    if not credentials:
        raise HTTPException(status_code=401, detail="인증 토큰이 누락되었습니다.")
    payload = decode_jwt_token(credentials.credentials)
    if not payload or "sub" not in payload:
        raise HTTPException(status_code=401, detail="유효하지 않거나 만료된 토큰입니다.")
    try:
        user_id = int(payload["sub"])
    except (ValueError, TypeError):
        raise HTTPException(status_code=401, detail="토큰의 sub 형식이 올바르지 않습니다.")
    user = db.get(User, user_id)
    if not user:
        raise HTTPException(status_code=401, detail="존재하지 않는 사용자입니다.")
    return user


def normalize_team_id(team_id: str | None = None) -> str:
    return (team_id or DEFAULT_TEAM_ID).strip() or DEFAULT_TEAM_ID


def folder_rule_db_id(space_id: str, folder_id: str) -> str:
    normalized_folder_id = str(folder_id or "other").strip() or "other"
    sid = (space_id or "").strip()
    return f"{sid}:{normalized_folder_id}" if sid else normalized_folder_id


def init_db() -> None:
    Base.metadata.create_all(bind=engine)
    ensure_document_schema()
    ensure_folder_rule_schema()
    ensure_review_queue_schema()
    db = SessionLocal()
    try:
        backfill_document_hashes(db)
        for team_id in TEAM_IDS:
            # 각 기본 팀에 기본 스페이스 확보(+폴더 규칙 시딩, 기존 문서 이관)
            ensure_team_default_space(db, team_id)
        db.commit()
    finally:
        db.close()


def ensure_document_schema() -> None:
    with engine.begin() as conn:
        columns = {row[1] for row in conn.execute(text("PRAGMA table_info(documents)")).fetchall()}
        if "file_hash" not in columns:
            conn.execute(text("ALTER TABLE documents ADD COLUMN file_hash VARCHAR(64)"))
            conn.execute(text("CREATE INDEX IF NOT EXISTS ix_documents_file_hash ON documents (file_hash)"))
        if "version_of_id" not in columns:
            conn.execute(text("ALTER TABLE documents ADD COLUMN version_of_id INTEGER"))
        if "team_id" not in columns:
            conn.execute(text("ALTER TABLE documents ADD COLUMN team_id VARCHAR(100) DEFAULT 'demo-team'"))
        if "space_id" not in columns:
            conn.execute(text("ALTER TABLE documents ADD COLUMN space_id VARCHAR(100)"))
        conn.execute(text("CREATE INDEX IF NOT EXISTS ix_documents_team_id ON documents (team_id)"))
        conn.execute(text("CREATE INDEX IF NOT EXISTS ix_documents_space_id ON documents (space_id)"))
        conn.execute(
            text("UPDATE documents SET team_id = :team_id WHERE team_id IS NULL OR team_id = ''"),
            {"team_id": DEFAULT_TEAM_ID},
        )


def ensure_folder_rule_schema() -> None:
    with engine.begin() as conn:
        columns = {row[1] for row in conn.execute(text("PRAGMA table_info(folder_rules)")).fetchall()}
        if "team_id" not in columns:
            conn.execute(text("ALTER TABLE folder_rules ADD COLUMN team_id VARCHAR(100) DEFAULT 'demo-team'"))
        if "folder_id" not in columns:
            conn.execute(text("ALTER TABLE folder_rules ADD COLUMN folder_id VARCHAR(100)"))
        if "icon" not in columns:
            conn.execute(text("ALTER TABLE folder_rules ADD COLUMN icon VARCHAR(100) DEFAULT ''"))
        if "reason" not in columns:
            conn.execute(text("ALTER TABLE folder_rules ADD COLUMN reason TEXT DEFAULT ''"))
        if "space_id" not in columns:
            conn.execute(text("ALTER TABLE folder_rules ADD COLUMN space_id VARCHAR(100)"))
        if "parent_folder_id" not in columns:
            conn.execute(text("ALTER TABLE folder_rules ADD COLUMN parent_folder_id VARCHAR(100)"))
        conn.execute(text("CREATE INDEX IF NOT EXISTS ix_folder_rules_team_id ON folder_rules (team_id)"))
        conn.execute(text("CREATE INDEX IF NOT EXISTS ix_folder_rules_folder_id ON folder_rules (folder_id)"))
        conn.execute(text("CREATE INDEX IF NOT EXISTS ix_folder_rules_space_id ON folder_rules (space_id)"))
        conn.execute(text("CREATE INDEX IF NOT EXISTS ix_folder_rules_parent_folder_id ON folder_rules (parent_folder_id)"))
        for rule in DEFAULT_FOLDER_RULES:
            conn.execute(
                text(
                    "UPDATE folder_rules "
                    "SET team_id = COALESCE(NULLIF(team_id, ''), :team_id), "
                    "folder_id = COALESCE(NULLIF(folder_id, ''), :folder_id), "
                    "icon = COALESCE(NULLIF(icon, ''), :icon), "
                    "reason = COALESCE(reason, '') "
                    "WHERE id = :legacy_id"
                ),
                {
                    "team_id": DEFAULT_TEAM_ID,
                    "folder_id": rule["id"],
                    "icon": FOLDER_ICON_MAP.get(rule["id"], ""),
                    "legacy_id": rule["id"],
                },
            )


def ensure_review_queue_schema() -> None:
    with engine.begin() as conn:
        columns = {row[1] for row in conn.execute(text("PRAGMA table_info(review_queue)")).fetchall()}
        if "team_id" not in columns:
            conn.execute(text("ALTER TABLE review_queue ADD COLUMN team_id VARCHAR(100) DEFAULT 'demo-team'"))
        if "suggested_category_id" not in columns:
            conn.execute(text("ALTER TABLE review_queue ADD COLUMN suggested_category_id VARCHAR(100) DEFAULT ''"))
        if "evidence_package" not in columns:
            conn.execute(text("ALTER TABLE review_queue ADD COLUMN evidence_package TEXT DEFAULT ''"))
        if "recommendations_json" not in columns:
            conn.execute(text("ALTER TABLE review_queue ADD COLUMN recommendations_json TEXT DEFAULT '[]'"))
        if "resolved_category" not in columns:
            conn.execute(text("ALTER TABLE review_queue ADD COLUMN resolved_category VARCHAR(255) DEFAULT ''"))
        if "resolved_category_id" not in columns:
            conn.execute(text("ALTER TABLE review_queue ADD COLUMN resolved_category_id VARCHAR(100) DEFAULT ''"))
        if "resolved_at" not in columns:
            conn.execute(text("ALTER TABLE review_queue ADD COLUMN resolved_at DATETIME"))
        if "updated_at" not in columns:
            conn.execute(text("ALTER TABLE review_queue ADD COLUMN updated_at DATETIME"))
        if "space_id" not in columns:
            conn.execute(text("ALTER TABLE review_queue ADD COLUMN space_id VARCHAR(100)"))
        conn.execute(text("CREATE INDEX IF NOT EXISTS ix_review_queue_team_id ON review_queue (team_id)"))
        conn.execute(text("CREATE INDEX IF NOT EXISTS ix_review_queue_status ON review_queue (status)"))
        conn.execute(text("CREATE INDEX IF NOT EXISTS ix_review_queue_space_id ON review_queue (space_id)"))
        conn.execute(
            text("UPDATE review_queue SET team_id = :team_id WHERE team_id IS NULL OR team_id = ''"),
            {"team_id": DEFAULT_TEAM_ID},
        )
        conn.execute(text("UPDATE review_queue SET status = 'pending' WHERE status IS NULL OR status = ''"))
        conn.execute(text("UPDATE review_queue SET reason = '' WHERE reason IS NULL"))
        conn.execute(text("UPDATE review_queue SET suggested_category = '' WHERE suggested_category IS NULL"))
        conn.execute(text("UPDATE review_queue SET suggested_category_id = '' WHERE suggested_category_id IS NULL"))
        conn.execute(text("UPDATE review_queue SET evidence_package = '' WHERE evidence_package IS NULL"))
        conn.execute(text("UPDATE review_queue SET recommendations_json = '[]' WHERE recommendations_json IS NULL OR recommendations_json = ''"))
        conn.execute(text("UPDATE review_queue SET resolved_category = '' WHERE resolved_category IS NULL"))
        conn.execute(text("UPDATE review_queue SET resolved_category_id = '' WHERE resolved_category_id IS NULL"))
        conn.execute(text("UPDATE review_queue SET updated_at = COALESCE(updated_at, created_at) WHERE updated_at IS NULL"))


def seed_folder_rules(db: Session, space_id: str, team_id: str, reset: bool = False) -> None:
    """스페이스 단위로 기본 분류 폴더 규칙을 시딩한다."""
    normalized_team_id = normalize_team_id(team_id)
    if reset:
        db.query(FolderRule).filter(FolderRule.space_id == space_id).delete(synchronize_session=False)
        db.flush()
    for order, rule in enumerate(DEFAULT_FOLDER_RULES):
        folder_id = rule["id"]
        existing = (
            db.query(FolderRule)
            .filter(FolderRule.space_id == space_id, FolderRule.folder_id == folder_id)
            .first()
        )
        payload = default_folder_rule_payload(rule, space_id, normalized_team_id)
        if existing:
            if reset:
                apply_folder_rule_payload(existing, payload)
            continue
        db.add(FolderRule(**payload))


def default_folder_rule_payload(rule: dict, space_id: str, team_id: str) -> dict:
    folder_id = rule["id"]
    return {
        "id": folder_rule_db_id(space_id, folder_id),
        "team_id": normalize_team_id(team_id),
        "space_id": space_id,
        "parent_folder_id": rule.get("parent_folder_id") or None,
        "folder_id": folder_id,
        "name": rule.get("name", folder_id),
        "icon": FOLDER_ICON_MAP.get(folder_id, ""),
        "description": rule.get("description", ""),
        "keywords": rule.get("keywords", ""),
        "context_terms": rule.get("context_terms", ""),
        "reason": rule.get("reason", ""),
        "updated_at": datetime.utcnow(),
    }


def create_default_space(db: Session, team_id: str, name: str = "기본 스페이스") -> Space:
    """팀에 스페이스를 생성하고 기본 분류 폴더 규칙을 시딩한다."""
    normalized_team_id = normalize_team_id(team_id)
    space_id = f"space-{uuid.uuid4().hex[:12]}"
    space = Space(id=space_id, team_id=normalized_team_id, name=name, description="")
    db.add(space)
    db.flush()
    seed_folder_rules(db, space_id, normalized_team_id)
    return space


def ensure_team_default_space(db: Session, team_id: str) -> Space:
    """팀의 첫 스페이스를 반환. 없으면 기본 스페이스를 만들고 기존 팀 문서를 이관한다."""
    normalized_team_id = normalize_team_id(team_id)
    space = (
        db.query(Space)
        .filter(Space.team_id == normalized_team_id)
        .order_by(Space.created_at)
        .first()
    )
    if space:
        return space
    space = create_default_space(db, normalized_team_id)
    # 기존 팀 문서/검토큐(스페이스 미지정)를 기본 스페이스로 이관
    db.query(Document).filter(
        Document.team_id == normalized_team_id, Document.space_id.is_(None)
    ).update({"space_id": space.id}, synchronize_session=False)
    db.query(ReviewQueue).filter(
        ReviewQueue.team_id == normalized_team_id, ReviewQueue.space_id.is_(None)
    ).update({"space_id": space.id}, synchronize_session=False)
    db.commit()
    return space


def resolve_space_id(db: Session, team_id: str, space_id: str | None) -> str:
    """team_id + 선택적 space_id를 받아 유효한 space_id로 해석한다 (없으면 기본 스페이스)."""
    normalized_team_id = normalize_team_id(team_id)
    sid = (space_id or "").strip()
    if sid:
        sp = db.query(Space).filter(Space.id == sid, Space.team_id == normalized_team_id).first()
        if sp:
            return sp.id
    return ensure_team_default_space(db, normalized_team_id).id


def apply_folder_rule_payload(rule: FolderRule, payload: dict) -> None:
    rule.team_id = payload["team_id"]
    if payload.get("space_id"):
        rule.space_id = payload["space_id"]
    if "parent_folder_id" in payload:
        rule.parent_folder_id = payload["parent_folder_id"]
    rule.folder_id = payload["folder_id"]
    rule.name = payload["name"]
    rule.icon = payload.get("icon", "")
    rule.description = payload.get("description", "")
    rule.keywords = payload.get("keywords", "")
    rule.context_terms = payload.get("context_terms", "")
    rule.reason = payload.get("reason", "")
    rule.updated_at = datetime.utcnow()


def backfill_document_hashes(db: Session) -> None:
    docs = db.query(Document).filter(Document.file_hash.is_(None)).all()
    changed = False
    for doc in docs:
        stored_path = Path(doc.stored_path or "")
        if not stored_path.exists():
            continue
        try:
            doc.file_hash = hashlib.sha256(stored_path.read_bytes()).hexdigest()
            changed = True
        except OSError:
            continue
    if changed:
        db.commit()


app = FastAPI(title="Dowple DMS PDF Analyzer", version="0.1.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)

init_db()


@app.get("/", response_class=HTMLResponse)
def index():
    html_path = BASE_DIR.parent / "dowple_dms.html"
    if html_path.exists():
        return HTMLResponse(content=html_path.read_text(encoding="utf-8"))
    raise HTTPException(status_code=404, detail="HTML file not found")


@app.get("/health")
def health() -> dict:
    return {"status": "ok"}


@app.post("/api/auth/google")
async def auth_google(payload: dict = Body(...), db: Session = Depends(get_db)) -> dict:
    credential = payload.get("credential")
    if not credential:
        raise HTTPException(status_code=400, detail="credential 토큰이 누락되었습니다.")

    email = "demo-user@example.com"
    name = "Demo User"
    picture = ""

    # 이메일 형식이면 가상 로그인 모드로 판별
    is_mock_login = False
    if "@" in credential and not credential.startswith("ey"):
        if re.match(r"[^@]+@[^@]+\.[^@]+", credential):
            is_mock_login = True

    try:
        if is_mock_login:
            # 가상 로그인 모드: 구글 OAuth 검증을 우회하고 계정 정보 자동 획득
            email = credential.strip()
            name = email.split("@")[0]
            picture = ""
        else:
            if GOOGLE_CLIENT_ID:
                # 실제 구글 ID 토큰 검증
                idinfo = id_token.verify_oauth2_token(credential, google_requests.Request(), GOOGLE_CLIENT_ID)
                email = idinfo["email"]
                name = idinfo.get("name", "")
                picture = idinfo.get("picture", "")
            else:
                # 개발/데모 시뮬레이션 모드 (클라이언트 ID 미설정 시)
                if "@" in credential:
                    email = credential.strip()
                    name = email.split("@")[0]
                else:
                    try:
                        # 전달받은 데이터가 JWT 형태인지 임시 확인 시도
                        decoded = jwt.decode(credential, options={"verify_signature": False})
                        email = decoded.get("email", email)
                        name = decoded.get("name", name)
                        picture = decoded.get("picture", picture)
                    except Exception:
                        pass
    except ValueError as e:
        raise HTTPException(status_code=400, detail=f"구글 토큰 검증 실패: {str(e)}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"인증 처리 오류: {str(e)}")

    user = db.query(User).filter(User.email == email).first()
    if not user:
        user = User(email=email, name=name, picture=picture)
        db.add(user)
        db.commit()
        db.refresh(user)
    else:
        # 프로필 정보 업데이트
        user.name = name
        user.picture = picture
        db.commit()
        db.refresh(user)

    token = create_jwt_token(user.id, user.email)
    return {
        "ok": True,
        "token": token,
        "user": {
            "id": user.id,
            "email": user.email,
            "name": user.name,
            "picture": user.picture
        }
    }


def check_team_access(team_id: str, user_id: int, db: Session):
    normalized = normalize_team_id(team_id)
    if normalized == "demo-team":
        return
    # 개인 팀: personal-{user_id} 형식이면 본인인지만 확인
    if normalized.startswith("personal-"):
        try:
            owner_id = int(normalized.split("-", 1)[1])
            if owner_id != user_id:
                raise HTTPException(status_code=403, detail="본인의 개인 팀에만 접근할 수 있습니다.")
        except (ValueError, IndexError):
            raise HTTPException(status_code=400, detail="올바르지 않은 개인 팀 ID입니다.")
        return
    membership = db.query(UserTeam).filter(UserTeam.user_id == user_id, UserTeam.team_id == normalized).first()
    if not membership:
        raise HTTPException(
            status_code=403,
            detail=f"해당 팀({normalized})에 접근할 권한이 없습니다. 먼저 팀에 가입해주세요."
        )


@app.get("/api/auth/me")
def get_auth_me(current_user: User = Depends(get_current_user)) -> dict:
    return {
        "ok": True,
        "user": {
            "id": current_user.id,
            "email": current_user.email,
            "name": current_user.name,
            "picture": current_user.picture
        }
    }


@app.post("/api/teams")
def create_team(payload: dict = Body(...), current_user: User = Depends(get_current_user), db: Session = Depends(get_db)) -> dict:
    team_name = str(payload.get("name") or "").strip()
    if not team_name:
        raise HTTPException(status_code=400, detail="팀 이름을 입력해주세요.")

    team_id = f"team-{uuid.uuid4().hex[:12]}"
    # 고유한 6자리 초대코드 생성
    invite_code = uuid.uuid4().hex[:6].upper()

    # 중복 초대코드 확인 및 재시도
    attempts = 0
    while db.query(Team).filter(Team.invite_code == invite_code).first() and attempts < 10:
        invite_code = uuid.uuid4().hex[:6].upper()
        attempts += 1

    team = Team(id=team_id, name=team_name, owner_id=current_user.id, invite_code=invite_code)
    db.add(team)
    db.flush()

    # 생성자를 소유자(owner) 역할로 가입
    user_team = UserTeam(user_id=current_user.id, team_id=team_id, role="owner")
    db.add(user_team)

    # 신규 팀에 기본 스페이스 생성 (스페이스 단위 기본 폴더 규칙도 함께 시딩됨)
    default_space = create_default_space(db, team_id)

    db.commit()
    db.refresh(team)

    return {
        "ok": True,
        "team": {
            "id": team.id,
            "name": team.name,
            "invite_code": team.invite_code,
            "role": "owner"
        }
    }


@app.post("/api/teams/join")
def join_team(payload: dict = Body(...), current_user: User = Depends(get_current_user), db: Session = Depends(get_db)) -> dict:
    invite_code = str(payload.get("invite_code") or "").strip().upper()
    if not invite_code:
        raise HTTPException(status_code=400, detail="초대 코드를 입력해주세요.")

    team = db.query(Team).filter(Team.invite_code == invite_code).first()
    if not team:
        raise HTTPException(status_code=404, detail="해당 초대 코드와 일치하는 팀이 존재하지 않습니다.")

    # 이미 가입되어 있는지 확인
    existing = db.query(UserTeam).filter(UserTeam.user_id == current_user.id, UserTeam.team_id == team.id).first()
    if existing:
        return {
            "ok": True,
            "message": "이미 가입된 팀입니다.",
            "team": {
                "id": team.id,
                "name": team.name,
                "role": existing.role
            }
        }

    # 신규 멤버로 추가
    user_team = UserTeam(user_id=current_user.id, team_id=team.id, role="member")
    db.add(user_team)
    db.commit()

    return {
        "ok": True,
        "team": {
            "id": team.id,
            "name": team.name,
            "role": "member"
        }
    }


@app.get("/api/teams")
def list_my_teams(current_user: User = Depends(get_current_user), db: Session = Depends(get_db)) -> dict:
    personal_space_id = f"personal-{current_user.id}"

    results = (
        db.query(Team, UserTeam.role)
        .join(UserTeam, Team.id == UserTeam.team_id)
        .filter(UserTeam.user_id == current_user.id)
        .all()
    )

    # 개인 팀을 항상 첫 번째 항목으로 추가
    teams_list = [{
        "id": personal_space_id,
        "name": "내 팀",
        "type": "personal",
        "role": "owner",
        "invite_code": ""
    }]
    for team, role in results:
        teams_list.append({
            "id": team.id,
            "name": team.name,
            "type": "team",
            "invite_code": team.invite_code,
            "role": role,
            "created_at": team.created_at.isoformat() if team.created_at else ""
        })

    return {
        "ok": True,
        "teams": teams_list
    }


@app.get("/api/teams/{team_id}/spaces")
def get_team_spaces(team_id: str, current_user: User = Depends(get_current_user), db: Session = Depends(get_db)) -> dict:
    """특정 팀의 스페이스 목록 조회 (스페이스가 없으면 기본 스페이스 자동 생성)"""
    normalized_team_id = normalize_team_id(team_id)
    check_team_access(normalized_team_id, current_user.id, db)

    # 스페이스가 없으면 기본 스페이스 확보(+폴더 규칙 시딩, 기존 문서 이관)
    ensure_team_default_space(db, normalized_team_id)
    spaces = (
        db.query(Space)
        .filter(Space.team_id == normalized_team_id)
        .order_by(Space.created_at)
        .all()
    )

    spaces_list = [
        {
            "id": space.id,
            "name": space.name,
            "description": space.description or "",
            "created_at": space.created_at.isoformat() if space.created_at else ""
        }
        for space in spaces
    ]
    return {
        "ok": True,
        "spaces": spaces_list
    }


@app.post("/api/spaces")
def create_space(payload: dict = Body(...), current_user: User = Depends(get_current_user), db: Session = Depends(get_db)) -> dict:
    """새로운 스페이스 생성

    folder_setup: 분류 폴더 초기 구성 방식
      - "default"(기본): 기본 16개 폴더 시딩
      - "copy": copy_from_space_id 스페이스의 폴더 규칙 복제
      - "empty": 시딩하지 않음 (구조 업로드 등으로 이후 채움)
    """
    team_id = str(payload.get("team_id") or "").strip()
    space_name = str(payload.get("name") or "").strip()
    folder_setup = str(payload.get("folder_setup") or "default").strip().lower()
    copy_from_space_id = str(payload.get("copy_from_space_id") or "").strip()

    if not team_id:
        raise HTTPException(status_code=400, detail="팀 ID를 입력해주세요.")
    if not space_name:
        raise HTTPException(status_code=400, detail="스페이스 이름을 입력해주세요.")

    check_team_access(team_id, current_user.id, db)
    normalized_team_id = normalize_team_id(team_id)

    space_id = f"space-{uuid.uuid4().hex[:12]}"
    space = Space(
        id=space_id,
        team_id=normalized_team_id,
        name=space_name,
        description=payload.get("description", "")
    )
    db.add(space)
    db.flush()

    if folder_setup == "copy" and copy_from_space_id:
        # 원본 스페이스의 폴더 규칙을 새 스페이스로 복제 (같은 팀 내에서만)
        source = db.query(Space).filter(
            Space.id == copy_from_space_id, Space.team_id == normalized_team_id
        ).first()
        if not source:
            raise HTTPException(status_code=400, detail="복사할 원본 스페이스를 찾을 수 없습니다.")
        source_rules = get_space_folder_rules(db, copy_from_space_id)
        if source_rules:
            for rule in source_rules:
                db.add(FolderRule(
                    id=folder_rule_db_id(space_id, rule.folder_id or rule.id),
                    team_id=normalized_team_id,
                    space_id=space_id,
                    parent_folder_id=rule.parent_folder_id,
                    folder_id=rule.folder_id,
                    name=rule.name,
                    icon=rule.icon,
                    description=rule.description,
                    keywords=rule.keywords,
                    context_terms=rule.context_terms,
                    reason=rule.reason,
                    updated_at=datetime.utcnow(),
                ))
        else:
            seed_folder_rules(db, space_id, normalized_team_id)
    elif folder_setup == "empty":
        pass  # 폴더 규칙 시딩하지 않음 (구조 업로드 등으로 이후 채움)
    else:
        # default: 기본 16개 폴더 시딩
        seed_folder_rules(db, space_id, normalized_team_id)

    db.commit()
    db.refresh(space)

    return {
        "ok": True,
        "space": {
            "id": space.id,
            "name": space.name,
            "description": space.description or "",
            "created_at": space.created_at.isoformat() if space.created_at else ""
        }
    }


@app.delete("/api/spaces/{space_id}")
def delete_space(space_id: str, current_user: User = Depends(get_current_user), db: Session = Depends(get_db)) -> dict:
    """스페이스 삭제"""
    space = db.query(Space).filter(Space.id == space_id).first()
    if not space:
        raise HTTPException(status_code=404, detail="스페이스를 찾을 수 없습니다.")

    check_team_access(space.team_id, current_user.id, db)

    db.query(Document).filter(Document.space_id == space_id).update({"space_id": None}, synchronize_session=False)
    db.query(FolderRule).filter(FolderRule.space_id == space_id).update({"space_id": None}, synchronize_session=False)
    db.query(ReviewQueue).filter(ReviewQueue.space_id == space_id).update({"space_id": None}, synchronize_session=False)
    db.delete(space)
    db.commit()

    return {"ok": True, "message": "스페이스가 삭제되었습니다."}


@app.post("/api/migrate-demo-to-personal")
def migrate_demo_to_personal(
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
) -> dict:
    """demo-team의 문서/폴더규칙/검토큐를 현재 사용자의 개인 팀으로 이전합니다."""
    personal_id = f"personal-{current_user.id}"

    # 이미 이전한 적 있는지 확인 (개인 팀에 이미 문서가 있으면 중복 이전 방지)
    already_has_personal = db.query(Document).filter(Document.team_id == personal_id).first()
    demo_doc_count = db.query(Document).filter(Document.team_id == "demo-team").count()

    migrated_docs = 0
    migrated_reviews = 0

    if demo_doc_count > 0:
        migrated_docs = db.query(Document).filter(Document.team_id == "demo-team").update(
            {"team_id": personal_id}, synchronize_session=False
        )
        migrated_reviews = db.query(ReviewQueue).filter(ReviewQueue.team_id == "demo-team").update(
            {"team_id": personal_id}, synchronize_session=False
        )

    # 개인 팀의 기본 스페이스 확보(+기본 폴더 규칙 시딩, 기존 문서 이관)
    ensure_team_default_space(db, personal_id)
    db.commit()

    return {
        "ok": True,
        "personal_space_id": personal_id,
        "migrated_docs": migrated_docs,
        "migrated_reviews": migrated_reviews
    }


@app.get("/api/folder-rules")
def list_folder_rules(
    team_id: str = Query(DEFAULT_TEAM_ID),
    space_id: str = Query(""),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
) -> dict:
    normalized_team_id = normalize_team_id(team_id)
    check_team_access(normalized_team_id, current_user.id, db)
    resolved_space_id = resolve_space_id(db, normalized_team_id, space_id)
    seed_folder_rules(db, resolved_space_id, normalized_team_id)
    db.commit()
    rules = get_space_folder_rules(db, resolved_space_id)
    return {
        "ok": True,
        "team_id": normalized_team_id,
        "space_id": resolved_space_id,
        "folder_rules": [serialize_folder_rule(rule) for rule in rules],
    }


@app.put("/api/folder-rules")
def update_folder_rules(
    payload: Any = Body(...),
    team_id: str = Query(DEFAULT_TEAM_ID),
    space_id: str = Query(""),
    replace: bool = Query(False),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
) -> dict:
    normalized_team_id = normalize_team_id(team_id)
    check_team_access(normalized_team_id, current_user.id, db)
    resolved_space_id = resolve_space_id(db, normalized_team_id, space_id)
    incoming_rules = payload.get("folder_rules", payload.get("rules", [])) if isinstance(payload, dict) else payload
    if not isinstance(incoming_rules, list):
        raise HTTPException(status_code=400, detail="folder_rules must be a list")

    # replace=True: 기존 폴더 규칙을 모두 비우고 들어온 규칙만 저장 (기본 16개 강제 추가 안 함)
    # → 폴더 구조 업로드로 "맞춤 폴더만" 구성할 때 사용
    if replace:
        db.query(FolderRule).filter(FolderRule.space_id == resolved_space_id).delete(synchronize_session=False)
        db.flush()
        for index, raw_rule in enumerate(incoming_rules):
            if not isinstance(raw_rule, dict):
                continue
            payload_rule = normalize_folder_rule_payload(raw_rule, resolved_space_id, normalized_team_id, index)
            db.add(FolderRule(**payload_rule))
        db.commit()
        rules = get_space_folder_rules(db, resolved_space_id)
        return {
            "ok": True,
            "team_id": normalized_team_id,
            "space_id": resolved_space_id,
            "folder_rules": [serialize_folder_rule(rule) for rule in rules],
        }

    existing_by_folder_id = {
        rule.folder_id or rule.id: rule
        for rule in get_space_folder_rules(db, resolved_space_id)
    }
    seen_folder_ids = set()
    for index, raw_rule in enumerate(incoming_rules):
        if not isinstance(raw_rule, dict):
            continue
        payload_rule = normalize_folder_rule_payload(raw_rule, resolved_space_id, normalized_team_id, index)
        folder_id = payload_rule["folder_id"]
        seen_folder_ids.add(folder_id)
        existing = existing_by_folder_id.get(folder_id)
        if existing:
            apply_folder_rule_payload(existing, payload_rule)
        else:
            db.add(FolderRule(**payload_rule))

    for default_rule in DEFAULT_FOLDER_RULES:
        if default_rule["id"] not in seen_folder_ids:
            payload_rule = default_folder_rule_payload(default_rule, resolved_space_id, normalized_team_id)
            existing = existing_by_folder_id.get(default_rule["id"])
            if existing:
                apply_folder_rule_payload(existing, payload_rule)
            else:
                db.add(FolderRule(**payload_rule))

    db.commit()
    rules = get_space_folder_rules(db, resolved_space_id)
    return {
        "ok": True,
        "team_id": normalized_team_id,
        "space_id": resolved_space_id,
        "folder_rules": [serialize_folder_rule(rule) for rule in rules],
    }


@app.post("/api/folder-rules/reset")
def reset_folder_rules_api(
    team_id: str = Query(DEFAULT_TEAM_ID),
    space_id: str = Query(""),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
) -> dict:
    normalized_team_id = normalize_team_id(team_id)
    check_team_access(normalized_team_id, current_user.id, db)
    resolved_space_id = resolve_space_id(db, normalized_team_id, space_id)
    seed_folder_rules(db, resolved_space_id, normalized_team_id, reset=True)
    db.commit()
    rules = get_space_folder_rules(db, resolved_space_id)
    return {
        "ok": True,
        "team_id": normalized_team_id,
        "space_id": resolved_space_id,
        "folder_rules": [serialize_folder_rule(rule) for rule in rules],
    }


@app.post("/api/folders")
def create_folder(payload: dict = Body(...), current_user: User = Depends(get_current_user), db: Session = Depends(get_db)) -> dict:
    """스페이스 안에 단일 폴더(루트 또는 하위)를 생성한다."""
    team_id = str(payload.get("team_id") or "").strip()
    space_id = str(payload.get("space_id") or "").strip()
    name = str(payload.get("name") or "").strip()
    parent_folder_id = str(payload.get("parent_folder_id") or "").strip() or None
    icon = str(payload.get("icon") or "ti-folder").strip() or "ti-folder"

    if not name:
        raise HTTPException(status_code=400, detail="폴더 이름을 입력해주세요.")
    normalized_team_id = normalize_team_id(team_id)
    check_team_access(normalized_team_id, current_user.id, db)
    resolved_space_id = resolve_space_id(db, normalized_team_id, space_id)

    # 부모 폴더 검증 (있으면 같은 스페이스에 존재해야 함)
    if parent_folder_id:
        parent = (
            db.query(FolderRule)
            .filter(FolderRule.space_id == resolved_space_id, FolderRule.folder_id == parent_folder_id)
            .first()
        )
        if not parent:
            raise HTTPException(status_code=400, detail="부모 폴더를 찾을 수 없습니다.")

    folder_id = f"folder-{uuid.uuid4().hex[:12]}"
    db.add(FolderRule(
        id=folder_rule_db_id(resolved_space_id, folder_id),
        team_id=normalized_team_id,
        space_id=resolved_space_id,
        parent_folder_id=parent_folder_id,
        folder_id=folder_id,
        name=name,
        icon=icon,
        description=str(payload.get("description") or ""),
        keywords=join_rule_terms(payload.get("keywords", "")),
        context_terms=join_rule_terms(payload.get("contextTerms", payload.get("context_terms", ""))),
        reason=str(payload.get("reason") or ""),
        updated_at=datetime.utcnow(),
    ))
    db.commit()
    rules = get_space_folder_rules(db, resolved_space_id)
    return {
        "ok": True,
        "team_id": normalized_team_id,
        "space_id": resolved_space_id,
        "folder": {"id": folder_id, "name": name, "parent_folder_id": parent_folder_id or ""},
        "folder_rules": [serialize_folder_rule(rule) for rule in rules],
    }


@app.delete("/api/folders/{folder_id}")
def delete_folder(folder_id: str, team_id: str = Query(DEFAULT_TEAM_ID), space_id: str = Query(""), current_user: User = Depends(get_current_user), db: Session = Depends(get_db)) -> dict:
    """폴더와 그 하위 폴더들을 삭제한다 (스페이스 범위). 문서는 미분류로 남긴다."""
    normalized_team_id = normalize_team_id(team_id)
    check_team_access(normalized_team_id, current_user.id, db)
    resolved_space_id = resolve_space_id(db, normalized_team_id, space_id)

    # 하위 폴더까지 재귀 수집
    all_rules = get_space_folder_rules(db, resolved_space_id)
    children_map: dict = {}
    for r in all_rules:
        children_map.setdefault(r.parent_folder_id or "", []).append(r.folder_id)
    to_delete = []
    stack = [folder_id]
    while stack:
        fid = stack.pop()
        to_delete.append(fid)
        stack.extend(children_map.get(fid, []))

    db.query(FolderRule).filter(
        FolderRule.space_id == resolved_space_id, FolderRule.folder_id.in_(to_delete)
    ).delete(synchronize_session=False)
    db.commit()
    rules = get_space_folder_rules(db, resolved_space_id)
    return {
        "ok": True,
        "deleted": to_delete,
        "folder_rules": [serialize_folder_rule(rule) for rule in rules],
    }


@app.patch("/api/folders/{folder_id}")
def rename_folder(folder_id: str, payload: dict = Body(...), team_id: str = Query(DEFAULT_TEAM_ID), space_id: str = Query(""), current_user: User = Depends(get_current_user), db: Session = Depends(get_db)) -> dict:
    """폴더 이름(및 아이콘)을 변경한다."""
    normalized_team_id = normalize_team_id(team_id)
    check_team_access(normalized_team_id, current_user.id, db)
    resolved_space_id = resolve_space_id(db, normalized_team_id, space_id)
    rule = (
        db.query(FolderRule)
        .filter(FolderRule.space_id == resolved_space_id, FolderRule.folder_id == folder_id)
        .first()
    )
    if not rule:
        raise HTTPException(status_code=404, detail="폴더를 찾을 수 없습니다.")
    new_name = str(payload.get("name") or "").strip()
    if new_name:
        rule.name = new_name
    if payload.get("icon"):
        rule.icon = str(payload.get("icon")).strip()
    rule.updated_at = datetime.utcnow()
    db.commit()
    rules = get_space_folder_rules(db, resolved_space_id)
    return {
        "ok": True,
        "folder": {"id": folder_id, "name": rule.name},
        "folder_rules": [serialize_folder_rule(rule) for rule in rules],
    }


@app.post("/api/analyze")
async def analyze(file: UploadFile = File(...), current_user: User = Depends(get_current_user)) -> dict:
    start = time.perf_counter()
    filename = file.filename or "uploaded_file"
    file_type = detect_file_type(filename, file.content_type or "")
    content = await file.read()

    if file_type in {"docx", "txt", "hwpx", "pptx", "xlsx", "csv", "hwp", "image"}:
        result = analyze_document_content(filename, file_type, content)
        return {
            "ok": True,
            "filename": filename,
            "file_type": file_type,
            "extraction_status": result["extraction_status"],
            "mode": result["mode"],
            "page_count": result["page_count"],
            "evidence_package": result["evidence_package"],
            "processing_time_sec": round(time.perf_counter() - start, 2),
            "warning": result["warning"],
        }

    if file_type != "pdf":
        return {
            "ok": True,
            "filename": filename,
            "file_type": file_type,
            "extraction_status": "unsupported",
            "mode": "fallback",
            "page_count": 0,
            "evidence_package": "",
            "processing_time_sec": round(time.perf_counter() - start, 2),
            "warning": "현재 PDF 외 형식은 시뮬레이션 처리",
        }

    try:
        evidence_package, page_count, total_text_len = extract_pdf_evidence(filename, content)
        if total_text_len < TEXT_LAYER_MIN_CHARS:
            return {
                "ok": True,
                "filename": filename,
                "file_type": "pdf",
                "extraction_status": "no_text_layer",
                "mode": "PDF 열기 성공 - 텍스트 레이어 없음",
                "page_count": page_count,
                "evidence_package": "",
                "processing_time_sec": round(time.perf_counter() - start, 2),
                "warning": "텍스트 레이어가 부족합니다. 스캔본 또는 이미지형 PDF일 수 있어 OCR 처리가 필요합니다.",
            }
        return {
            "ok": True,
            "filename": filename,
            "file_type": "pdf",
            "extraction_status": "success",
            "mode": "PyMuPDF 실제 텍스트 추출",
            "page_count": page_count,
            "evidence_package": evidence_package,
            "processing_time_sec": round(time.perf_counter() - start, 2),
            "warning": "",
        }
    except Exception:
        return {
            "ok": False,
            "filename": filename,
            "file_type": "pdf",
            "extraction_status": "failed",
            "mode": "fallback",
            "page_count": 0,
            "evidence_package": "",
            "processing_time_sec": 0,
            "warning": "PDF 텍스트 추출 실패",
        }


@app.post("/api/documents")
async def create_document(
    file: UploadFile = File(...),
    team_id: str = Form(DEFAULT_TEAM_ID),
    space_id: str = Form(""),
    category: str = Form(""),
    category_id: str = Form(""),
    summary: str = Form(""),
    confidence: float = Form(0),
    client_evidence_package: str = Form(""),
    duplicate_decision: str = Form(""),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
) -> dict:
    normalized_team_id = normalize_team_id(team_id)
    check_team_access(normalized_team_id, current_user.id, db)

    # space_id 검증 (없거나 유효하지 않으면 팀의 기본 스페이스로 귀속)
    if space_id:
        space = db.query(Space).filter(Space.id == space_id, Space.team_id == normalized_team_id).first()
        if not space:
            raise HTTPException(status_code=400, detail="해당 팀에 존재하는 스페이스가 아닙니다.")
    else:
        space_id = ensure_team_default_space(db, normalized_team_id).id
    filename = file.filename or "uploaded_file"
    file_type = detect_file_type(filename, file.content_type or "")
    content = await file.read()
    file_hash = hashlib.sha256(content).hexdigest()
    evidence_package = ""
    extraction_status = "unsupported"
    page_count = 0
    warning = ""

    if file_type in {"docx", "txt", "hwpx", "pptx", "xlsx", "csv", "hwp", "image"}:
        result = analyze_document_content(filename, file_type, content)
        evidence_package = result["evidence_package"]
        extraction_status = result["extraction_status"]
        page_count = result["page_count"]
        warning = result["warning"]
    elif file_type == "pdf":
        try:
            evidence_package, page_count, total_text_len = extract_pdf_evidence(filename, content)
            if total_text_len < TEXT_LAYER_MIN_CHARS:
                evidence_package = ""
                extraction_status = "no_text_layer"
                warning = "텍스트 레이어가 부족합니다. 스캔본 또는 이미지형 PDF일 수 있어 OCR 처리가 필요합니다."
            else:
                extraction_status = "success"
        except Exception:
            extraction_status = "failed"
            warning = "PDF 텍스트 추출 실패"

    if not evidence_package and client_evidence_package:
        evidence_package = truncate_text(client_evidence_package, MAX_EVIDENCE_CHARS)
        # 클라이언트(브라우저 OCR 등)에서 텍스트를 확보한 경우, 추출 실패 상태를 성공으로 갱신
        if extraction_status in {"ocr_required", "no_text_layer", "unsupported", "failed"}:
            extraction_status = "success"
            warning = ""

    duplicate_candidates = find_duplicate_candidates(
        db=db,
        team_id=normalized_team_id,
        filename=filename,
        file_hash=file_hash,
        evidence_package=evidence_package,
    )
    duplicate_type = duplicate_candidates[0]["duplicate_type"] if duplicate_candidates else ""
    normalized_decision = (duplicate_decision or "").strip().lower()

    if normalized_decision == "check_only":
        return {
            "ok": True,
            "document": None,
            "team_id": normalized_team_id,
            "duplicate_type": duplicate_type,
            "duplicate_candidates": duplicate_candidates,
            "file_hash": file_hash,
            "extraction_status": extraction_status,
            "page_count": page_count,
            "warning": warning,
        }

    version_of_id = None
    if normalized_decision in {"version", "save_as_version"} and duplicate_candidates:
        version_of_id = duplicate_candidates[0]["id"]

    stored_path = save_uploaded_file(filename, content)

    doc = Document(
        team_id=normalized_team_id,
        space_id=space_id if space_id else None,
        filename=filename,
        file_type=file_type,
        category=category,
        category_id=category_id,
        summary=summary,
        confidence=confidence,
        evidence_package=evidence_package,
        extraction_status=extraction_status,
        page_count=page_count,
        file_hash=file_hash,
        version_of_id=version_of_id,
        stored_path=str(stored_path),
        archived=False,
    )
    db.add(doc)
    db.commit()
    db.refresh(doc)

    if confidence < 40 or category_id == "other":
        existing_review = (
            db.query(ReviewQueue)
            .filter(
                ReviewQueue.team_id == normalized_team_id,
                ReviewQueue.filename == filename,
                ReviewQueue.status == "pending",
            )
            .first()
        )
        if existing_review:
            existing_review.document_id = doc.id
            existing_review.suggested_category = category
            existing_review.suggested_category_id = category_id
            existing_review.confidence = confidence
            existing_review.evidence_package = evidence_package
            existing_review.updated_at = datetime.utcnow()
        else:
            db.add(ReviewQueue(
                team_id=normalized_team_id,
                space_id=space_id if space_id else None,
                document_id=doc.id,
                filename=filename,
                reason="신뢰도 낮음 또는 기타 분류",
                suggested_category=category,
                suggested_category_id=category_id,
                confidence=confidence,
                evidence_package=evidence_package,
                recommendations_json="[]",
                status="pending",
                updated_at=datetime.utcnow(),
            ))
        db.commit()

    return {
        "ok": True,
        "document": serialize_document(doc),
        "team_id": normalized_team_id,
        "duplicate_type": duplicate_type,
        "duplicate_candidates": duplicate_candidates,
        "warning": warning,
    }


@app.get("/api/documents")
def list_documents(
    team_id: str = Query(DEFAULT_TEAM_ID),
    space_id: str = Query(""),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
) -> dict:
    normalized_team_id = normalize_team_id(team_id)
    check_team_access(normalized_team_id, current_user.id, db)
    query = (
        db.query(Document)
        .filter(Document.archived == False, Document.team_id == normalized_team_id)  # noqa: E712
    )
    if space_id:
        query = query.filter(Document.space_id == space_id)
    docs = query.order_by(Document.created_at.desc()).all()
    return {"ok": True, "documents": [serialize_document(doc) for doc in docs]}


@app.get("/api/documents/{document_id}")
def get_document(
    document_id: int,
    team_id: str | None = Query(None),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
) -> dict:
    doc = db.get(Document, document_id)
    if not doc or doc.archived:
        raise HTTPException(status_code=404, detail="Document not found")
    check_team_access(doc.team_id, current_user.id, db)
    if team_id is not None and doc.team_id != normalize_team_id(team_id):
        raise HTTPException(status_code=404, detail="Document not found")
    return {"ok": True, "document": serialize_document(doc)}


@app.patch("/api/documents/{document_id}")
def update_document(
    document_id: int,
    payload: Any = Body(...),
    team_id: str = Query(DEFAULT_TEAM_ID),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
) -> dict:
    normalized_team_id = normalize_team_id(team_id)
    check_team_access(normalized_team_id, current_user.id, db)
    doc = db.get(Document, document_id)
    if not doc or doc.archived:
        raise HTTPException(status_code=404, detail="Document not found")
    if doc.team_id != normalized_team_id:
        raise HTTPException(status_code=404, detail="Document not found")
    if not isinstance(payload, dict):
        raise HTTPException(status_code=400, detail="JSON body is required")

    if "category" in payload:
        doc.category = str(payload.get("category") or "")
    if "category_id" in payload:
        doc.category_id = str(payload.get("category_id") or "")
    if "summary" in payload:
        doc.summary = str(payload.get("summary") or "")
    if "confidence" in payload:
        try:
            doc.confidence = float(payload.get("confidence") or 0)
        except (TypeError, ValueError):
            raise HTTPException(status_code=400, detail="confidence must be a number") from None
    if "evidence_package" in payload:
        doc.evidence_package = truncate_text(str(payload.get("evidence_package") or ""), MAX_EVIDENCE_CHARS)
    if "extraction_status" in payload:
        doc.extraction_status = str(payload.get("extraction_status") or doc.extraction_status or "unsupported")
    if "page_count" in payload:
        try:
            doc.page_count = int(payload.get("page_count") or 0)
        except (TypeError, ValueError):
            raise HTTPException(status_code=400, detail="page_count must be an integer") from None

    doc.updated_at = datetime.utcnow()
    db.commit()
    db.refresh(doc)
    return {"ok": True, "document": serialize_document(doc)}


@app.get("/api/documents/{document_id}/download")
def download_document(
    document_id: int,
    team_id: str | None = Query(None),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
) -> FileResponse:
    doc = db.get(Document, document_id)
    if not doc or doc.archived:
        raise HTTPException(status_code=404, detail="Document not found")
    check_team_access(doc.team_id, current_user.id, db)
    if team_id is not None and doc.team_id != normalize_team_id(team_id):
        raise HTTPException(status_code=404, detail="Document not found")

    stored_path = Path(doc.stored_path or "")
    if not stored_path.exists() or not stored_path.is_file():
        raise HTTPException(status_code=404, detail="Stored file not found")

    media_type = "application/pdf" if (doc.file_type == "pdf" or doc.filename.lower().endswith(".pdf")) else "application/octet-stream"
    return FileResponse(
        path=stored_path,
        media_type=media_type,
        filename=doc.filename,
        content_disposition_type="inline",
    )


@app.delete("/api/documents/{document_id}")
def delete_document(
    document_id: int,
    team_id: str = Query(DEFAULT_TEAM_ID),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
) -> dict:
    normalized_team_id = normalize_team_id(team_id)
    check_team_access(normalized_team_id, current_user.id, db)
    doc = db.get(Document, document_id)
    if not doc or doc.team_id != normalized_team_id:
        raise HTTPException(status_code=404, detail="Document not found")
    doc.archived = True
    doc.updated_at = datetime.utcnow()
    db.commit()
    db.refresh(doc)
    return {"ok": True, "document": serialize_document(doc)}


@app.get("/api/review-queue")
def list_review_queue(
    team_id: str = Query(DEFAULT_TEAM_ID),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
) -> dict:
    normalized_team_id = normalize_team_id(team_id)
    check_team_access(normalized_team_id, current_user.id, db)
    items = (
        db.query(ReviewQueue)
        .filter(ReviewQueue.team_id == normalized_team_id, ReviewQueue.status == "pending")
        .order_by(ReviewQueue.created_at.desc(), ReviewQueue.id.desc())
        .all()
    )
    return {"ok": True, "review_queue": [serialize_review_queue_item(item) for item in items]}


@app.post("/api/review-queue")
def create_review_queue_item(
    payload: Any = Body(...),
    team_id: str = Query(DEFAULT_TEAM_ID),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
) -> dict:
    normalized_team_id = normalize_team_id(team_id)
    check_team_access(normalized_team_id, current_user.id, db)
    if not isinstance(payload, dict):
        raise HTTPException(status_code=400, detail="JSON body is required")

    filename = str(payload.get("filename") or "").strip()
    if not filename:
        raise HTTPException(status_code=400, detail="filename is required")

    existing = (
        db.query(ReviewQueue)
        .filter(
            ReviewQueue.team_id == normalized_team_id,
            ReviewQueue.filename == filename,
            ReviewQueue.status == "pending",
        )
        .order_by(ReviewQueue.created_at.desc(), ReviewQueue.id.desc())
        .first()
    )
    if existing:
        return {"ok": True, "review": serialize_review_queue_item(existing), "duplicate": True}

    document_id = payload.get("document_id")
    if document_id in ("", None):
        document_id = None
    else:
        try:
            document_id = int(document_id)
        except (TypeError, ValueError):
            raise HTTPException(status_code=400, detail="document_id must be an integer") from None

    try:
        confidence = float(payload.get("confidence") or 0)
    except (TypeError, ValueError):
        raise HTTPException(status_code=400, detail="confidence must be a number") from None

    now = datetime.utcnow()
    item = ReviewQueue(
        team_id=normalized_team_id,
        document_id=document_id,
        filename=filename,
        reason=str(payload.get("reason") or ""),
        suggested_category=str(payload.get("suggested_category") or ""),
        suggested_category_id=str(payload.get("suggested_category_id") or ""),
        confidence=confidence,
        evidence_package=truncate_text(str(payload.get("evidence_package") or ""), MAX_EVIDENCE_CHARS),
        recommendations_json=normalize_recommendations_json(payload.get("recommendations_json")),
        status="pending",
        created_at=now,
        updated_at=now,
    )
    db.add(item)
    db.commit()
    db.refresh(item)
    return {"ok": True, "review": serialize_review_queue_item(item), "duplicate": False}


@app.patch("/api/review-queue/{review_id}")
def update_review_queue_item(
    review_id: int,
    payload: Any = Body(...),
    team_id: str = Query(DEFAULT_TEAM_ID),
    db: Session = Depends(get_db),
) -> dict:
    normalized_team_id = normalize_team_id(team_id)
    if not isinstance(payload, dict):
        raise HTTPException(status_code=400, detail="JSON body is required")

    item = db.get(ReviewQueue, review_id)
    if not item or normalize_team_id(item.team_id) != normalized_team_id:
        raise HTTPException(status_code=404, detail="Review queue item not found")

    if "status" in payload:
        status = str(payload.get("status") or "pending").strip().lower()
        if status not in {"pending", "resolved", "ignored"}:
            raise HTTPException(status_code=400, detail="status must be pending, resolved, or ignored")
        item.status = status
        if status == "resolved":
            item.resolved_at = datetime.utcnow()

    if "resolved_category" in payload:
        item.resolved_category = str(payload.get("resolved_category") or "")
    if "resolved_category_id" in payload:
        item.resolved_category_id = str(payload.get("resolved_category_id") or "")
    if "document_id" in payload:
        document_id = payload.get("document_id")
        if document_id in ("", None):
            item.document_id = None
        else:
            try:
                item.document_id = int(document_id)
            except (TypeError, ValueError):
                raise HTTPException(status_code=400, detail="document_id must be an integer") from None

    item.updated_at = datetime.utcnow()
    db.commit()
    db.refresh(item)
    return {"ok": True, "review": serialize_review_queue_item(item)}


@app.delete("/api/review-queue/{review_id}")
def ignore_review_queue_item(
    review_id: int,
    team_id: str = Query(DEFAULT_TEAM_ID),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
) -> dict:
    normalized_team_id = normalize_team_id(team_id)
    check_team_access(normalized_team_id, current_user.id, db)
    item = db.get(ReviewQueue, review_id)
    if not item or normalize_team_id(item.team_id) != normalized_team_id:
        raise HTTPException(status_code=404, detail="Review queue item not found")
    item.status = "ignored"
    item.updated_at = datetime.utcnow()
    db.commit()
    db.refresh(item)
    return {"ok": True, "review": serialize_review_queue_item(item)}


def find_duplicate_candidates(db: Session, team_id: str, filename: str, file_hash: str, evidence_package: str) -> list[dict]:
    docs = (
        db.query(Document)
        .filter(Document.archived == False, Document.team_id == normalize_team_id(team_id))  # noqa: E712
        .order_by(Document.created_at.desc())
        .all()
    )
    candidates = []
    current_name_tokens = similarity_tokens(filename, is_filename=True)
    current_evidence_tokens = similarity_tokens(evidence_package)

    for doc in docs:
        if doc.file_hash and doc.file_hash == file_hash:
            candidates.append(build_duplicate_candidate(doc, 100, "exact_duplicate"))
            continue

        name_score = token_similarity(current_name_tokens, similarity_tokens(doc.filename, is_filename=True))
        evidence_score = token_similarity(current_evidence_tokens, similarity_tokens(doc.evidence_package or doc.summary or ""))
        similarity_score = round(max(name_score, evidence_score))
        if similarity_score >= DUPLICATE_VERSION_THRESHOLD:
            candidates.append(build_duplicate_candidate(doc, similarity_score, "possible_version"))

    candidates.sort(key=lambda item: (item["similarity_score"], item["created_at"]), reverse=True)
    return candidates[:5]


def build_duplicate_candidate(doc: Document, similarity_score: int, duplicate_type: str) -> dict:
    return {
        "id": doc.id,
        "team_id": doc.team_id or DEFAULT_TEAM_ID,
        "filename": doc.filename,
        "category": doc.category,
        "confidence": doc.confidence,
        "similarity_score": clamp_int(similarity_score, 0, 100),
        "duplicate_type": duplicate_type,
        "created_at": doc.created_at.isoformat() if doc.created_at else "",
    }


def similarity_tokens(text_value: str, is_filename: bool = False) -> list[str]:
    source = str(text_value or "").lower()
    if is_filename:
        source = re.sub(r"\.[a-z0-9]{1,10}$", " ", source)
        source = re.sub(r"(^|[\s._-])(v\d+)(?=$|[\s._-])", " ", source)
        source = re.sub(r"수정|최종|진짜최종|찐최종|final|copy|복사본", " ", source)

    raw_tokens = re.findall(r"[가-힣A-Za-z0-9]{2,}", source)
    stop_words = {
        "pdf", "docx", "hwp", "hwpx", "png", "jpg", "jpeg", "file", "document",
        "the", "and", "for", "with", "from", "this", "that", "문서", "자료",
    }
    tokens = []
    for token in raw_tokens:
        if token.isdigit() or token in stop_words:
            continue
        tokens.append(token)
    return list(dict.fromkeys(tokens[:120]))


def token_similarity(left: list[str], right: list[str]) -> float:
    if not left or not right:
        return 0
    left_set = set(left)
    right_set = set(right)
    intersection = left_set & right_set
    union = left_set | right_set
    return (len(intersection) / len(union)) * 100 if union else 0


def clamp_int(value: int, minimum: int, maximum: int) -> int:
    return max(minimum, min(maximum, int(value)))


def detect_file_type(filename: str, content_type: str) -> str:
    lower = filename.lower()
    if lower.endswith(".pdf") or content_type == "application/pdf":
        return "pdf"
    if lower.endswith(".hwpx"):
        return "hwpx"
    if lower.endswith(".hwp"):
        return "hwp"
    if lower.endswith(".docx"):
        return "docx"
    if lower.endswith(".pptx") or content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return "pptx"
    if lower.endswith(".ppt"):
        return "ppt"
    if lower.endswith(".xlsx") or content_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        return "xlsx"
    if lower.endswith(".xls"):
        return "xls"
    if lower.endswith(".csv") or content_type.startswith("text/csv") or content_type == "application/csv":
        return "csv"
    if lower.endswith(".txt") or content_type.startswith("text/"):
        return "txt"
    if lower.endswith((".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff")):
        return "image"
    return "document"


def save_uploaded_file(filename: str, content: bytes) -> Path:
    suffix = Path(filename).suffix.lower() or ".bin"
    safe_suffix = suffix if re.fullmatch(r"\.[A-Za-z0-9]{1,10}", suffix) else ".bin"
    stored_name = f"{uuid.uuid4().hex}{safe_suffix}"
    stored_path = UPLOAD_DIR / stored_name
    stored_path.write_bytes(content)
    return stored_path


def serialize_document(doc: Document) -> dict:
    return {
        "id": doc.id,
        "team_id": doc.team_id or DEFAULT_TEAM_ID,
        "space_id": doc.space_id or "",
        "filename": doc.filename,
        "file_type": doc.file_type,
        "category": doc.category,
        "category_id": doc.category_id,
        "summary": doc.summary,
        "confidence": doc.confidence,
        "evidence_package": doc.evidence_package,
        "extraction_status": doc.extraction_status,
        "page_count": doc.page_count,
        "file_hash": doc.file_hash,
        "version_of_id": doc.version_of_id,
        "stored_path": doc.stored_path,
        "created_at": doc.created_at.isoformat() if doc.created_at else "",
        "updated_at": doc.updated_at.isoformat() if doc.updated_at else "",
        "archived": doc.archived,
    }


def normalize_recommendations_json(value: Any) -> str:
    if isinstance(value, str):
        stripped = value.strip()
        if not stripped:
            return "[]"
        try:
            parsed = json.loads(stripped)
            return json.dumps(parsed if isinstance(parsed, list) else [parsed], ensure_ascii=False)
        except json.JSONDecodeError:
            return json.dumps([stripped], ensure_ascii=False)
    if isinstance(value, list):
        return json.dumps(value, ensure_ascii=False)
    if isinstance(value, dict):
        return json.dumps([value], ensure_ascii=False)
    return "[]"


def parse_recommendations_json(value: str | None) -> list:
    if not value:
        return []
    try:
        parsed = json.loads(value)
    except json.JSONDecodeError:
        return []
    return parsed if isinstance(parsed, list) else [parsed]


def serialize_review_queue_item(item: ReviewQueue) -> dict:
    return {
        "id": item.id,
        "team_id": item.team_id or DEFAULT_TEAM_ID,
        "document_id": item.document_id,
        "filename": item.filename,
        "reason": item.reason or "",
        "suggested_category": item.suggested_category or "",
        "suggested_category_id": item.suggested_category_id or "",
        "confidence": item.confidence or 0,
        "evidence_package": item.evidence_package or "",
        "recommendations": parse_recommendations_json(item.recommendations_json),
        "status": item.status or "pending",
        "resolved_category": item.resolved_category or "",
        "resolved_category_id": item.resolved_category_id or "",
        "created_at": item.created_at.isoformat() if item.created_at else "",
        "updated_at": item.updated_at.isoformat() if item.updated_at else "",
        "resolved_at": item.resolved_at.isoformat() if item.resolved_at else "",
    }


def get_team_folder_rules(db: Session, team_id: str) -> list[FolderRule]:
    normalized_team_id = normalize_team_id(team_id)
    rules = (
        db.query(FolderRule)
        .filter(FolderRule.team_id == normalized_team_id)
        .all()
    )
    order = {rule["id"]: index for index, rule in enumerate(DEFAULT_FOLDER_RULES)}
    return sorted(rules, key=lambda rule: order.get(rule.folder_id or rule.id, 999))


def get_space_folder_rules(db: Session, space_id: str) -> list[FolderRule]:
    rules = (
        db.query(FolderRule)
        .filter(FolderRule.space_id == space_id)
        .all()
    )
    order = {rule["id"]: index for index, rule in enumerate(DEFAULT_FOLDER_RULES)}
    return sorted(rules, key=lambda rule: order.get(rule.folder_id or rule.id, 999))


def serialize_folder_rule(rule: FolderRule) -> dict:
    folder_id = rule.folder_id or rule.id
    return {
        "id": folder_id,
        "folder_id": folder_id,
        "team_id": rule.team_id or DEFAULT_TEAM_ID,
        "space_id": rule.space_id or "",
        "parent_folder_id": rule.parent_folder_id or "",
        "name": rule.name,
        "icon": rule.icon or FOLDER_ICON_MAP.get(folder_id, ""),
        "description": rule.description or "",
        "keywords": split_rule_terms(rule.keywords),
        "contextTerms": split_rule_terms(rule.context_terms),
        "context_terms": split_rule_terms(rule.context_terms),
        "reason": rule.reason or "",
    }


def normalize_folder_rule_payload(raw_rule: dict, space_id: str, team_id: str, index: int = 0) -> dict:
    base = DEFAULT_FOLDER_RULES[index] if index < len(DEFAULT_FOLDER_RULES) else DEFAULT_FOLDER_RULES[-1]
    folder_id = str(raw_rule.get("folder_id") or raw_rule.get("id") or base["id"]).strip() or base["id"]
    parent_fid = str(raw_rule.get("parent_folder_id") or raw_rule.get("parentId") or "").strip() or None
    return {
        "id": folder_rule_db_id(space_id, folder_id),
        "team_id": normalize_team_id(team_id),
        "space_id": space_id,
        "parent_folder_id": parent_fid,
        "folder_id": folder_id,
        "name": str(raw_rule.get("name") or base.get("name") or folder_id).strip(),
        "icon": str(raw_rule.get("icon") or FOLDER_ICON_MAP.get(folder_id, "")).strip(),
        "description": str(raw_rule.get("description") or "").strip(),
        "keywords": join_rule_terms(raw_rule.get("keywords", base.get("keywords", ""))),
        "context_terms": join_rule_terms(raw_rule.get("contextTerms", raw_rule.get("context_terms", base.get("context_terms", "")))),
        "reason": str(raw_rule.get("reason") or "").strip(),
        "updated_at": datetime.utcnow(),
    }


def split_rule_terms(value: str | None) -> list[str]:
    return [term.strip() for term in str(value or "").replace("\n", ",").split(",") if term.strip()]


def join_rule_terms(value: Any) -> str:
    if isinstance(value, list):
        return ", ".join(str(item).strip() for item in value if str(item).strip())
    return ", ".join(split_rule_terms(str(value or "")))


def analyze_document_content(filename: str, file_type: str, content: bytes) -> dict:
    handlers = {
        "pdf": extract_pdf_evidence_result,
        "docx": extract_docx_evidence,
        "txt": extract_txt_evidence,
        "hwpx": extract_hwpx_evidence,
        "pptx": extract_pptx_evidence,
        "xlsx": extract_xlsx_evidence,
        "csv": extract_csv_evidence,
    }

    if file_type == "hwp":
        return {
            "extraction_status": "unsupported",
            "mode": "fallback",
            "page_count": 0,
            "evidence_package": "",
            "warning": "HWP 바이너리 파일은 향후 전용 파서 또는 변환 모듈이 필요합니다.",
        }
    if file_type == "image":
        return {
            "extraction_status": "ocr_required",
            "mode": "OCR 필요",
            "page_count": 0,
            "evidence_package": "",
            "warning": "이미지 문서는 OCR 처리가 필요합니다.",
        }

    handler = handlers.get(file_type)
    if not handler:
        return {
            "extraction_status": "unsupported",
            "mode": "fallback",
            "page_count": 0,
            "evidence_package": "",
            "warning": "현재 지원하지 않는 파일 형식입니다.",
        }
    return handler(filename, content)


def extract_pdf_evidence_result(filename: str, content: bytes) -> dict:
    evidence_package, page_count, total_text_len = extract_pdf_evidence(filename, content)
    if total_text_len < TEXT_LAYER_MIN_CHARS:
        return {
            "extraction_status": "no_text_layer",
            "mode": "PDF 열기 성공 - 텍스트 레이어 없음",
            "page_count": page_count,
            "evidence_package": "",
            "warning": "텍스트 레이어가 부족합니다. 스캔본 또는 이미지형 PDF일 수 있어 OCR 처리가 필요합니다.",
        }
    return {
        "extraction_status": "success",
        "mode": "PyMuPDF 실제 텍스트 추출",
        "page_count": page_count,
        "evidence_package": evidence_package,
        "warning": "",
    }


def build_text_evidence_package(filename: str, text_value: str, source_label: str, page_count: int = 0, count_label: str = "총 페이지 수") -> str:
    normalized = normalize_whitespace(text_value)
    first_text = truncate_text(normalized[:1800], 1800)
    last_text = truncate_text(normalized[-1200:], 1200) if len(normalized) > 1200 else normalized
    keyword_sentences = find_keyword_sentences(normalized)
    top_words = top_frequent_terms(normalized)
    key_sentences = extract_key_sentence_candidates(normalized)

    sections = [
        f"파일명: {filename}",
        f"추출 방식: {source_label}",
        f"전체 텍스트 길이: {len(normalized)}",
    ]
    if page_count:
        sections.append(f"{count_label}: {page_count}")
    sections.extend([
        "[첫 부분 텍스트]",
        first_text or "추출된 텍스트가 없습니다.",
        "[마지막 부분 텍스트]",
        last_text or "추출된 텍스트가 없습니다.",
        "[분류 키워드 포함 문장 최대 30개]",
        "\n".join(f"- {sentence}" for sentence in keyword_sentences) or "분류 키워드가 포함된 문장을 찾지 못했습니다.",
        "[자주 등장하는 단어 Top 20]",
        ", ".join(f"{word}({count})" for word, count in top_words) or "집계할 단어가 부족합니다.",
        "[핵심 문장 후보 최대 10개]",
        "\n".join(f"- {sentence}" for sentence in key_sentences) or "핵심 문장 후보가 부족합니다.",
    ])
    return truncate_text("\n\n".join(sections), MAX_EVIDENCE_CHARS)


def extract_docx_evidence(filename: str, content: bytes) -> dict:
    try:
        doc = DocxDocument(BytesIO(content))
        chunks = [paragraph.text for paragraph in doc.paragraphs if paragraph.text and paragraph.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
                if cells:
                    chunks.append(" | ".join(cells))
        text_value = "\n".join(chunks)
        if len(normalize_whitespace(text_value)) < 1:
            raise ValueError("DOCX text is empty")
        return {
            "extraction_status": "success",
            "mode": "python-docx 실제 텍스트 추출",
            "page_count": 0,
            "evidence_package": build_text_evidence_package(filename, text_value, "python-docx 실제 텍스트 추출"),
            "warning": "",
        }
    except Exception:
        return {
            "extraction_status": "failed",
            "mode": "fallback",
            "page_count": 0,
            "evidence_package": "",
            "warning": "DOCX 텍스트 추출 실패",
        }


def extract_txt_evidence(filename: str, content: bytes) -> dict:
    for encoding in ("utf-8-sig", "utf-8", "cp949", "euc-kr"):
        try:
            text_value = content.decode(encoding)
            if len(normalize_whitespace(text_value)) < 1:
                raise ValueError("TXT text is empty")
            return {
                "extraction_status": "success",
                "mode": "TXT 실제 텍스트 추출",
                "page_count": 0,
                "evidence_package": build_text_evidence_package(filename, text_value, f"TXT 실제 텍스트 추출 ({encoding})"),
                "warning": "",
            }
        except UnicodeDecodeError:
            continue
        except Exception:
            break
    return {
        "extraction_status": "failed",
        "mode": "fallback",
        "page_count": 0,
        "evidence_package": "",
        "warning": "TXT 텍스트 디코딩 실패",
    }


def extract_hwpx_evidence(filename: str, content: bytes) -> dict:
    try:
        chunks: list[str] = []
        with zipfile.ZipFile(BytesIO(content)) as archive:
            names = archive.namelist()
            candidates = [
                name for name in names
                if name.lower() == "preview/prvtext.txt"
                or (name.lower().startswith(("contents/", "sections/")) and name.lower().endswith(".xml"))
            ]
            for name in candidates:
                data = archive.read(name)
                if name.lower().endswith(".txt"):
                    chunks.append(decode_text_bytes(data))
                else:
                    chunks.append(extract_xml_text(data))
        text_value = "\n".join(chunk for chunk in chunks if chunk and chunk.strip())
        if len(normalize_whitespace(text_value)) < TEXT_LAYER_MIN_CHARS:
            raise ValueError("HWPX text is too short")
        return {
            "extraction_status": "success",
            "mode": "HWPX XML 실제 텍스트 추출",
            "page_count": 0,
            "evidence_package": build_text_evidence_package(filename, text_value, "HWPX XML 실제 텍스트 추출"),
            "warning": "",
        }
    except Exception:
        return {
            "extraction_status": "failed",
            "mode": "fallback",
            "page_count": 0,
            "evidence_package": "",
            "warning": "HWPX XML 텍스트 추출 실패",
        }


def extract_pptx_evidence(filename: str, content: bytes) -> dict:
    try:
        presentation = Presentation(BytesIO(content))
        slide_chunks: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            chunks: list[str] = []
            for shape in slide.shapes:
                chunks.extend(extract_pptx_shape_text(shape))
            try:
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        chunks.append(f"발표자 노트: {notes_text}")
            except Exception:
                pass

            unique_chunks = list(dict.fromkeys(chunk.strip() for chunk in chunks if chunk and chunk.strip()))
            if unique_chunks:
                slide_chunks.append(f"[슬라이드 {index}]\n" + "\n".join(unique_chunks))

        text_value = "\n\n".join(slide_chunks)
        slide_count = len(presentation.slides)
        if len(normalize_whitespace(text_value)) < TEXT_LAYER_MIN_CHARS:
            return {
                "extraction_status": "no_text_layer",
                "mode": "PPTX 열기 성공 - 텍스트 부족",
                "page_count": slide_count,
                "evidence_package": "",
                "warning": "PPTX에서 추출 가능한 텍스트가 부족합니다. 이미지형 슬라이드는 OCR 처리가 필요합니다.",
            }
        return {
            "extraction_status": "success",
            "mode": "python-pptx 실제 텍스트 추출",
            "page_count": slide_count,
            "evidence_package": build_text_evidence_package(
                filename,
                text_value,
                "python-pptx 실제 텍스트 추출",
                page_count=slide_count,
                count_label="총 슬라이드 수",
            ),
            "warning": "",
        }
    except Exception:
        return {
            "extraction_status": "failed",
            "mode": "fallback",
            "page_count": 0,
            "evidence_package": "",
            "warning": "PPTX 텍스트 추출 실패",
        }


def extract_pptx_shape_text(shape) -> list[str]:
    chunks: list[str] = []

    try:
        if getattr(shape, "has_text_frame", False) and shape.text_frame:
            paragraph_texts = []
            for paragraph in shape.text_frame.paragraphs:
                run_text = "".join(run.text for run in paragraph.runs).strip()
                paragraph_text = run_text or paragraph.text.strip()
                if paragraph_text:
                    paragraph_texts.append(paragraph_text)
            if paragraph_texts:
                chunks.append("\n".join(paragraph_texts))
    except Exception:
        pass

    try:
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
                if cells:
                    chunks.append(" | ".join(cells))
    except Exception:
        pass

    try:
        if hasattr(shape, "shapes"):
            for child_shape in shape.shapes:
                chunks.extend(extract_pptx_shape_text(child_shape))
    except Exception:
        pass

    try:
        if getattr(shape, "has_chart", False):
            chart = shape.chart
            if getattr(chart, "has_title", False) and chart.chart_title and chart.chart_title.has_text_frame:
                title_text = chart.chart_title.text_frame.text.strip()
                if title_text:
                    chunks.append(f"차트 제목: {title_text}")
    except Exception:
        pass

    return chunks


def extract_xlsx_evidence(filename: str, content: bytes) -> dict:
    try:
        workbook = load_workbook(BytesIO(content), read_only=True, data_only=True)
        sheet_names = list(workbook.sheetnames)
        rows: list[str] = []
        headers: list[str] = []

        try:
            for worksheet in workbook.worksheets:
                sheet_rows: list[str] = []
                for row_index, row in enumerate(worksheet.iter_rows(max_row=50, max_col=20, values_only=True), start=1):
                    cells = [safe_cell_text(value) for value in row]
                    cells = [cell for cell in cells if cell]
                    if not cells:
                        continue
                    row_text = f"{worksheet.title} R{row_index}: " + " | ".join(cells)
                    sheet_rows.append(row_text)
                    if row_index <= 3 and len(headers) < 20:
                        headers.append(f"{worksheet.title}: " + " | ".join(cells))
                rows.extend(sheet_rows)
        finally:
            workbook.close()

        if len(normalize_whitespace(" ".join(rows))) < TEXT_LAYER_MIN_CHARS:
            return {
                "extraction_status": "no_text_layer",
                "mode": "XLSX 열기 성공 - 표 텍스트 부족",
                "page_count": len(sheet_names),
                "evidence_package": "",
                "warning": "표에서 추출 가능한 텍스트가 부족합니다.",
            }

        return {
            "extraction_status": "success",
            "mode": "openpyxl 실제 XLSX 텍스트 추출",
            "page_count": len(sheet_names),
            "evidence_package": build_table_evidence_package(
                filename=filename,
                source_label="openpyxl 실제 XLSX 텍스트 추출",
                rows=rows,
                headers=headers,
                sheet_names=sheet_names,
                sheet_count=len(sheet_names),
            ),
            "warning": "",
        }
    except Exception:
        return {
            "extraction_status": "failed",
            "mode": "fallback",
            "page_count": 0,
            "evidence_package": "",
            "warning": "XLSX 텍스트 추출 실패",
        }


def extract_csv_evidence(filename: str, content: bytes) -> dict:
    for encoding in ("utf-8-sig", "utf-8", "cp949", "euc-kr"):
        try:
            text_value = content.decode(encoding)
        except UnicodeDecodeError:
            continue

        try:
            sample = text_value[:4096]
            try:
                dialect = csv.Sniffer().sniff(sample)
            except csv.Error:
                dialect = csv.excel

            reader = csv.reader(StringIO(text_value), dialect)
            rows: list[str] = []
            header = ""
            for row_index, row in enumerate(reader, start=1):
                if row_index > 100:
                    break
                cells = [safe_cell_text(value) for value in row[:30]]
                cells = [cell for cell in cells if cell]
                if not cells:
                    continue
                row_text = f"Row {row_index}: " + " | ".join(cells)
                rows.append(row_text)
                if not header:
                    header = " | ".join(cells)

            if len(normalize_whitespace(" ".join(rows))) < TEXT_LAYER_MIN_CHARS:
                return {
                    "extraction_status": "no_text_layer",
                    "mode": "CSV 열기 성공 - 표 텍스트 부족",
                    "page_count": 0,
                    "evidence_package": "",
                    "warning": "표에서 추출 가능한 텍스트가 부족합니다.",
                }

            return {
                "extraction_status": "success",
                "mode": "CSV 실제 텍스트 추출",
                "page_count": 0,
                "evidence_package": build_table_evidence_package(
                    filename=filename,
                    source_label="CSV 실제 텍스트 추출",
                    rows=rows,
                    headers=[header] if header else [],
                    encoding=encoding,
                ),
                "warning": "",
            }
        except Exception:
            break

    return {
        "extraction_status": "failed",
        "mode": "fallback",
        "page_count": 0,
        "evidence_package": "",
        "warning": "CSV 텍스트 추출 실패",
    }


def build_table_evidence_package(
    filename: str,
    source_label: str,
    rows: list[str],
    headers: list[str] | None = None,
    sheet_names: list[str] | None = None,
    encoding: str = "",
    sheet_count: int = 0,
) -> str:
    headers = headers or []
    sheet_names = sheet_names or []
    normalized_rows = [normalize_whitespace(row) for row in rows if normalize_whitespace(row)]
    normalized_text = "\n".join(normalized_rows)
    keyword_rows = find_keyword_sentences(normalized_text)
    top_words = top_frequent_terms(normalized_text)
    key_rows = extract_key_sentence_candidates(normalized_text)

    sections = [
        f"파일명: {filename}",
        f"추출 방식: {source_label}",
        f"전체 텍스트 길이: {len(normalize_whitespace(normalized_text))}",
    ]
    if sheet_count:
        sections.append(f"전체 시트 수: {sheet_count}")
    if sheet_names:
        sections.append(f"시트명 목록: {', '.join(sheet_names[:30])}")
    if encoding:
        sections.append(f"사용 인코딩: {encoding}")
    sections.extend([
        "[헤더 후보]",
        "\n".join(f"- {truncate_text(header, 260)}" for header in headers[:20]) or "헤더 후보가 부족합니다.",
        "[상위 행 일부]",
        "\n".join(f"- {truncate_text(row, 260)}" for row in normalized_rows[:35]) or "추출된 행 텍스트가 부족합니다.",
        "[분류 키워드 포함 행/문장 최대 30개]",
        "\n".join(f"- {row}" for row in keyword_rows[:30]) or "분류 키워드가 포함된 행/문장을 찾지 못했습니다.",
        "[자주 등장하는 단어 Top 20]",
        ", ".join(f"{word}({count})" for word, count in top_words) or "집계할 단어가 부족합니다.",
        "[핵심 행 후보 최대 10개]",
        "\n".join(f"- {row}" for row in key_rows[:10]) or "핵심 행 후보가 부족합니다.",
    ])
    return truncate_text("\n\n".join(sections), MAX_EVIDENCE_CHARS)


def safe_cell_text(value: Any, max_chars: int = 120) -> str:
    if value is None:
        return ""
    text_value = normalize_whitespace(str(value))
    if not text_value:
        return ""
    return truncate_text(text_value, max_chars)


def decode_text_bytes(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "cp949", "euc-kr"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="ignore")


def extract_xml_text(data: bytes) -> str:
    try:
        root = ElementTree.fromstring(data)
        return " ".join(text.strip() for text in root.itertext() if text and text.strip())
    except Exception:
        raw = decode_text_bytes(data)
        return normalize_whitespace(re.sub(r"<[^>]+>", " ", raw))


def extract_pdf_evidence(filename: str, content: bytes) -> tuple[str, int, int]:
    with fitz.open(stream=content, filetype="pdf") as doc:
        page_count = doc.page_count
        page_texts = [normalize_whitespace(page.get_text("text")) for page in doc]

    all_text = "\n".join(page_texts)
    total_text_len = len(normalize_whitespace(all_text))
    if total_text_len < TEXT_LAYER_MIN_CHARS:
        return "", page_count, total_text_len
    return build_text_evidence_package(filename, all_text, "PyMuPDF 실제 텍스트 추출", page_count=page_count), page_count, total_text_len

    first_pages = "\n\n".join(text for text in page_texts[:2] if text)
    last_page = page_texts[-1] if page_texts else ""
    keyword_sentences = find_keyword_sentences(all_text)
    top_words = top_frequent_words(all_text)

    sections = [
        f"파일명: {filename}",
        f"PDF 총 페이지 수: {page_count}",
        "[첫 1~2페이지 텍스트]",
        first_pages or "추출된 텍스트가 없습니다.",
        "[마지막 페이지 텍스트]",
        last_page or "추출된 텍스트가 없습니다.",
        "[분류 키워드 포함 문장 최대 30개]",
        "\n".join(f"- {sentence}" for sentence in keyword_sentences) or "분류 키워드가 포함된 문장을 찾지 못했습니다.",
        "[자주 등장하는 단어 Top 20]",
        ", ".join(f"{word}({count})" for word, count in top_words) or "집계할 단어가 부족합니다.",
    ]
    return truncate_text("\n\n".join(sections), MAX_EVIDENCE_CHARS), page_count, total_text_len


def find_keyword_sentences(text: str) -> List[str]:
    normalized = normalize_whitespace(text)
    sentences = re.split(r"(?<=[.!?。！？])\s+|\n+", normalized)
    found: List[str] = []
    seen = set()

    for sentence in sentences:
        clean = sentence.strip()
        if len(clean) < 8:
            continue
        if any(keyword.lower() in clean.lower() for keyword in KEYWORDS) or any(term.lower() in clean.lower() for term in KEY_SENTENCE_PRIORITY_TERMS):
            shortened = truncate_text(clean, 220)
            key = shortened.lower()
            if key not in seen:
                seen.add(key)
                found.append(shortened)
        if len(found) >= MAX_KEYWORD_SENTENCES:
            break

    return found


def extract_key_sentence_candidates(text: str) -> List[str]:
    normalized = normalize_whitespace(text)
    sentences = re.split(r"(?<=[.!?。！？])\s+|\n+|(?<=다\.)\s+", normalized)
    scored: list[tuple[int, str]] = []
    seen = set()
    date_pattern = re.compile(r"\d{4}[./-]\d{1,2}[./-]\d{1,2}|\d{1,2}[./-]\d{1,2}")
    money_pattern = re.compile(r"\d[\d,]*(원|만원|천원|억원|KRW|USD|달러)", re.IGNORECASE)
    email_pattern = re.compile(r"[\w.+-]+@[\w-]+(?:\.[\w-]+)+")
    phone_pattern = re.compile(r"\d{2,4}[-.\s]\d{3,4}[-.\s]\d{4}")
    org_pattern = re.compile(r"[가-힣A-Za-z0-9]+(?:기관|대학교|협회|재단|센터|연구원|주식회사|회사|Inc|Corp|LLC)")

    for sentence in sentences:
        clean = sentence.strip()
        if len(clean) < 18:
            continue
        shortened = truncate_text(clean, 240)
        key = shortened.lower()
        if key in seen:
            continue
        seen.add(key)
        lower = clean.lower()
        score = 0
        if any(term.lower() in lower for term in KEY_SENTENCE_PRIORITY_TERMS):
            score += 4
        if date_pattern.search(clean):
            score += 3
        if money_pattern.search(clean):
            score += 3
        if email_pattern.search(clean):
            score += 2
        if phone_pattern.search(clean):
            score += 2
        if org_pattern.search(clean):
            score += 2
        if score:
            scored.append((score, shortened))

    scored.sort(key=lambda item: (-item[0], len(item[1])))
    return [sentence for _, sentence in scored[:10]]


def top_frequent_terms(text: str) -> List[tuple[str, int]]:
    tokens = re.findall(r"[가-힣A-Za-z0-9]{2,}", text.lower())
    filtered = [
        token for token in tokens
        if token not in STOP_WORDS and not token.isdigit()
    ]
    return Counter(filtered).most_common(20)


def top_frequent_words(text: str) -> List[tuple[str, int]]:
    tokens = re.findall(r"[가-힣A-Za-z0-9]{2,}", text.lower())
    filtered = [
        token for token in tokens
        if token not in STOP_WORDS and not token.isdigit()
    ]
    return Counter(filtered).most_common(20)


def normalize_whitespace(text: str) -> str:
    return re.sub(r"\s+", " ", text or "").strip()


def truncate_text(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 20].rstrip() + "\n... [truncated]"
